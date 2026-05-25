# -*- coding: utf-8 -*-
"""
Jubiaba Smart City - Apresentacao para Investidores 2025/2026
Geracao automatica de PPTX a partir da Estrutura de Slides v4 (47 slides).

Paleta KV oficial (definida na estrutura): Preto + Branco + Turquesa
Tipografia: Calibri (sistema) com fallback EI brand (Cormorant / DM Sans).
"""

import os
import glob
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ----------------------------------------------------------------------
# CONFIG
# ----------------------------------------------------------------------
OUT = r"D:\Copy EI\jubiaba\Jubiaba_Apresentacao_Investidores_2025.pptx"
PHOTOS_AREA = r"D:\Copy EI\jubiaba\Investidores SERGIPE\FOTOS_\Fotos da Área"
PHOTOS_REGIAO = r"D:\Copy EI\jubiaba\Investidores SERGIPE\FOTOS_\Fotos da Região_"
PHOTO_TOTAL = r"D:\Copy EI\jubiaba\Investidores SERGIPE\ÁREA TOTAL JUBIABÁ.jpeg"
PHOTO_JUR = r"D:\Copy EI\jubiaba\Investidores SERGIPE\Simulação de Jurerê, sobreposta em roxo pelo tamanho do terreno de Jubiabá.png"

# Brand palette - Layout KV
VOID = RGBColor(0x0A, 0x0A, 0x0A)        # Preto profundo
NAVY = RGBColor(0x0C, 0x18, 0x2A)        # Card navy
NAVY_LIGHT = RGBColor(0x14, 0x24, 0x3D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xF0, 0xED, 0xE6)
TURQ = RGBColor(0x00, 0xC2, 0xB3)        # Turquesa KV
TURQ_DARK = RGBColor(0x00, 0x8F, 0x86)
GOLD = RGBColor(0xB8, 0x94, 0x2E)        # Acento (de EI brand)
GREY_50 = RGBColor(0x9A, 0x9A, 0x9A)
GREY_30 = RGBColor(0x55, 0x55, 0x55)
GREY_15 = RGBColor(0x22, 0x22, 0x22)

# Fonts (Calibri = onipresente Windows; serif p/ destaques)
FONT_SERIF = "Cambria"   # serif robusta com fallback universal
FONT_SANS = "Calibri"
FONT_MONO = "Consolas"

# Slide dim 16:9 widescreen 33.87 x 19.05 cm
SLIDE_W = Cm(33.867)
SLIDE_H = Cm(19.05)

# ----------------------------------------------------------------------
# Helpers
# ----------------------------------------------------------------------

def find_image(folder, keywords=None, fallback_index=0):
    """Encontra a primeira imagem em folder cujo nome contem alguma keyword."""
    if not os.path.isdir(folder):
        return None
    files = sorted(
        f for f in os.listdir(folder)
        if f.lower().endswith(('.jpg', '.jpeg', '.png', '.webp'))
    )
    if not files:
        return None
    if keywords:
        for kw in keywords:
            for f in files:
                if kw.lower() in f.lower():
                    return os.path.join(folder, f)
    return os.path.join(folder, files[fallback_index % len(files)])


def list_images(folder, limit=8):
    if not os.path.isdir(folder):
        return []
    files = sorted(
        f for f in os.listdir(folder)
        if f.lower().endswith(('.jpg', '.jpeg', '.png', '.webp'))
    )
    return [os.path.join(folder, f) for f in files[:limit]]


def add_bg(slide, color):
    """Pinta fundo do slide com retangulo full-bleed."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    # Send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_image_full(slide, path):
    if path and os.path.exists(path):
        try:
            slide.shapes.add_picture(path, 0, 0, width=SLIDE_W, height=SLIDE_H)
            return True
        except Exception:
            return False
    return False


def add_overlay(slide, opacity=0.55, color=VOID):
    """Camada escura traduzida em retangulo semi-transparente."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    rect.line.fill.background()
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    # Aplicar transparencia via XML
    sp = rect.fill._xPr
    solidFill = sp.find(qn('a:solidFill'))
    if solidFill is not None:
        srgb = solidFill.find(qn('a:srgbClr'))
        if srgb is not None:
            alpha = etree.SubElement(srgb, qn('a:alpha'))
            alpha.set('val', str(int((1 - opacity) * 100000)))
    return rect


def add_text(slide, x, y, w, h, text, *, font=FONT_SANS, size=14,
             color=CREAM, bold=False, italic=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.15, letter_spacing=0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = text.split('\n') if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        f = run.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
        if letter_spacing:
            # spc em centesimos de ponto
            rPr = run._r.get_or_add_rPr()
            rPr.set('spc', str(int(letter_spacing * 100)))
    return tb


def add_rect(slide, x, y, w, h, fill=NAVY, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    r = slide.shapes.add_shape(shape_type, x, y, w, h)
    if radius:
        r.adjustments[0] = 0.06
    if fill is None:
        r.fill.background()
    else:
        r.fill.solid()
        r.fill.fore_color.rgb = fill
    if line is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line
        r.line.width = Pt(0.75)
    r.shadow.inherit = False
    return r


def add_line(slide, x1, y1, x2, y2, color=TURQ, weight=1.5):
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def add_card(slide, x, y, w, h, title, body, *, accent=TURQ,
             title_size=14, body_size=10):
    add_rect(slide, x, y, w, h, fill=NAVY, line=None, radius=True)
    # Accent bar lateral
    bar = add_rect(slide, x, y, Cm(0.12), h, fill=accent)
    bar.line.fill.background()
    add_text(slide, x + Cm(0.5), y + Cm(0.35), w - Cm(0.7), Cm(1.1),
             title, font=FONT_SERIF, size=title_size, color=WHITE, bold=True)
    add_text(slide, x + Cm(0.5), y + Cm(1.4), w - Cm(0.7), h - Cm(1.6),
             body, font=FONT_SANS, size=body_size, color=CREAM,
             line_spacing=1.25)


def add_section_header(slide, eyebrow, title, *, color_eye=TURQ,
                       color_title=WHITE, top=Cm(1.3)):
    # Eyebrow
    add_text(slide, Cm(1.5), top, Cm(20), Cm(0.7), eyebrow,
             font=FONT_SANS, size=10, color=color_eye, bold=True,
             letter_spacing=4)
    # Hairline
    add_line(slide, Cm(1.5), top + Cm(0.85), Cm(4.5), top + Cm(0.85),
             color=color_eye, weight=1.2)
    # Title
    add_text(slide, Cm(1.5), top + Cm(1.0), Cm(30), Cm(2.6), title,
             font=FONT_SERIF, size=34, color=color_title, bold=False,
             line_spacing=1.05)


def add_footer(slide, n_slide, total=47, color=GREY_50, scheme="dark"):
    bg_text = WHITE if scheme == "dark" else VOID
    # Logo lockup
    add_text(slide, Cm(1.5), SLIDE_H - Cm(1.0), Cm(15), Cm(0.6),
             "GSWP  ·  KV  ·  HAMAM GLOBAL",
             font=FONT_SANS, size=8, color=color, bold=True,
             letter_spacing=4)
    # Slide number
    add_text(slide, SLIDE_W - Cm(4), SLIDE_H - Cm(1.0), Cm(2.5), Cm(0.6),
             f"{n_slide:02d} / {total:02d}",
             font=FONT_SANS, size=8, color=color, align=PP_ALIGN.RIGHT,
             letter_spacing=2)


def add_disclaimer(slide):
    add_text(slide, Cm(1.5), SLIDE_H - Cm(0.55), SLIDE_W - Cm(3), Cm(0.4),
             "Projecoes baseadas em analise de mercado e historico de "
             "valorizacao. Retornos passados nao garantem resultados futuros. "
             "Consulte o prospecto completo antes de investir.",
             font=FONT_SANS, size=6, color=GREY_30, italic=True)


# ----------------------------------------------------------------------
# Build presentation
# ----------------------------------------------------------------------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def new_slide(bg_color=VOID):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, bg_color)
    return slide

# ----------------------------------------------------------------------
# Pre-cuts: imagens curadas
# ----------------------------------------------------------------------
img_aerea = find_image(PHOTOS_AREA, ['13.14.07', '13.14.05', '13.14.10'])
img_aerea2 = find_image(PHOTOS_AREA, ['13.14.10 (2)', '13.14.10 (3)'], 1)
img_aerea3 = find_image(PHOTOS_AREA, ['13.14.09'], 2)
img_aero = find_image(PHOTOS_AREA, ['AEROPORTO'])
img_praia_saco = find_image(PHOTOS_REGIAO, ['praia-do-saco-um-paraiso'])
img_mangue = find_image(PHOTOS_REGIAO, ['mangue-seco'])
img_lagoa = find_image(PHOTOS_REGIAO, ['lagoa-dos-tambaquis'])
img_dunas = find_image(PHOTOS_REGIAO, ['dunas'])
img_porsol = find_image(PHOTOS_REGIAO, ['por-do-sol'])
img_nascer = find_image(PHOTOS_REGIAO, ['nascer-do-sol'])
img_aerea_extra = find_image(PHOTOS_REGIAO, ['dji-1599'])

gallery_area = list_images(PHOTOS_AREA, 6)
gallery_regiao = [
    img_praia_saco, img_mangue, img_lagoa, img_dunas, img_porsol, img_nascer
]

# ======================================================================
# BLOCO 01 - ABERTURA & ENCANTAMENTO
# ======================================================================

# --- S01 CAPA --------------------------------------------------------
s = new_slide(VOID)
if img_aerea:
    add_image_full(s, img_aerea)
add_overlay(s, opacity=0.55, color=VOID)
# Sutil gradiente inferior
grad = add_rect(s, 0, Cm(13), SLIDE_W, Cm(6.05), fill=VOID)
sp = grad.fill._xPr.find(qn('a:solidFill'))
srgb = sp.find(qn('a:srgbClr'))
alpha = etree.SubElement(srgb, qn('a:alpha'))
alpha.set('val', '40000')

# Top eyebrow
add_text(s, Cm(1.5), Cm(1.2), Cm(20), Cm(0.6),
         "APRESENTACAO PARA INVESTIDORES  ·  2025 / 2026",
         font=FONT_SANS, size=9, color=TURQ, bold=True, letter_spacing=6)
# Hairline
add_line(s, Cm(1.5), Cm(1.95), Cm(7), Cm(1.95), color=TURQ, weight=1.2)

# Title
add_text(s, Cm(1.5), Cm(7.5), Cm(28), Cm(4.5),
         "JUBIABA",
         font=FONT_SERIF, size=120, color=WHITE, bold=False,
         letter_spacing=20, line_spacing=1.0)
add_text(s, Cm(1.5), Cm(12.0), Cm(28), Cm(1.0),
         "A PROXIMA FRONTEIRA DO LUXO NO NORDESTE",
         font=FONT_SANS, size=12, color=CREAM, letter_spacing=8)

# Bottom partners
add_line(s, Cm(1.5), Cm(16.4), Cm(32.3), Cm(16.4), color=TURQ, weight=0.75)
add_text(s, Cm(1.5), Cm(16.7), Cm(31), Cm(0.8),
         "GSWP    ×    KONIGSBERGER VANNUCCHI    ×    HAMAM GLOBAL",
         font=FONT_SANS, size=10, color=WHITE, bold=True,
         letter_spacing=8, align=PP_ALIGN.CENTER)
add_text(s, Cm(1.5), Cm(17.5), Cm(31), Cm(0.6),
         "Praia dos Abais  ·  Estancia  ·  Sergipe  ·  Brasil",
         font=FONT_SERIF, size=11, color=GREY_50, italic=True,
         align=PP_ALIGN.CENTER, letter_spacing=4)


# --- S02 O MOMENTO ---------------------------------------------------
s = new_slide(VOID)
add_section_header(s, "01  ·  O MOMENTO",
                   "O litoral virgem como ativo.\nA tecnologia integrada desde a fundacao.")

# Manifesto (esquerda)
manifesto = (
    "Um projeto raro nasce na confluencia de tres fatores irrepetiveis: "
    "uma janela regulatoria que se abriu apos uma decada de impasse, "
    "uma faixa de litoral virgem com 1.070m de praia exclusiva, "
    "e uma triade executora de classe mundial.\n\n"
    "Jubiaba nao e um retrofit. Nao e correcao de um master plan antigo. "
    "E uma smart city de planta nova, projetada com tecnologia subterranea "
    "desde a primeira linha do desenho.\n\n"
    "A janela de entrada como socio fundador esta aberta. "
    "Por tempo limitado."
)
add_text(s, Cm(1.5), Cm(6.5), Cm(17.5), Cm(11), manifesto,
         font=FONT_SERIF, size=15, color=CREAM, line_spacing=1.4)

# Checklist (direita)
checklist_x = Cm(20.5)
add_rect(s, checklist_x, Cm(6.5), Cm(11.8), Cm(11), fill=NAVY, radius=True)
add_text(s, checklist_x + Cm(0.6), Cm(6.9), Cm(11), Cm(0.8),
         "OS SEIS PILARES",
         font=FONT_SANS, size=9, color=TURQ, bold=True, letter_spacing=4)

itens = [
    "Nao e retrofit",
    "Litoral virgem como ativo",
    "Tecnologia integrada desde a fundacao",
    "Paisagem como protagonista",
    "Desenvolvimento sustentavel",
    "Smart city de planta nova",
]
for i, it in enumerate(itens):
    yy = Cm(8.0 + i * 1.45)
    # Marker
    mk = add_rect(s, checklist_x + Cm(0.6), yy + Cm(0.25), Cm(0.45), Cm(0.45),
                  fill=TURQ, radius=True)
    add_text(s, checklist_x + Cm(0.6), yy + Cm(0.18), Cm(0.45), Cm(0.6),
             "+", font=FONT_SERIF, size=14, color=VOID, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, checklist_x + Cm(1.4), yy + Cm(0.2), Cm(10), Cm(0.7),
             it, font=FONT_SERIF, size=14, color=WHITE)

add_footer(s, 2)


# --- S03 SUMARIO EXECUTIVO ------------------------------------------
s = new_slide(VOID)
add_section_header(s, "02  ·  SUMARIO EXECUTIVO",
                   "Seis numeros que definem o projeto.")

stats = [
    ("129", "hectares", "Fase 1   ·   potencial +70ha"),
    ("1.070m", "de praia", "Faixa exclusiva pe na areia"),
    ("R$ 2bi", "VGV", "Projetado total do complexo"),
    ("3", "parceiros", "GSWP  ·  KV  ·  Hamam Global"),
    ("R$ 10mi", "Round 02", "Captacao em andamento"),
    ("Fase 1", "agora", "Maxima janela como socio"),
]
gx = Cm(1.5)
gy = Cm(6.5)
gw = Cm(10.05)
gh = Cm(5.4)
for i, (num, label, sub) in enumerate(stats):
    col = i % 3
    row = i // 3
    x = gx + col * (gw + Cm(0.4))
    y = gy + row * (gh + Cm(0.4))
    add_rect(s, x, y, gw, gh, fill=NAVY, radius=True)
    add_rect(s, x, y, Cm(0.12), gh, fill=TURQ)
    add_text(s, x + Cm(0.7), y + Cm(0.5), gw - Cm(0.9), Cm(2.4),
             num, font=FONT_SERIF, size=46, color=WHITE,
             line_spacing=1.0)
    add_text(s, x + Cm(0.7), y + Cm(2.9), gw - Cm(0.9), Cm(0.7),
             label.upper(), font=FONT_SANS, size=10, color=TURQ,
             bold=True, letter_spacing=4)
    add_text(s, x + Cm(0.7), y + Cm(3.7), gw - Cm(0.9), Cm(1.5),
             sub, font=FONT_SANS, size=10, color=CREAM, line_spacing=1.3)

add_footer(s, 3)


# --- S04 A TRIADE EXECUTORA -----------------------------------------
s = new_slide(VOID)
add_section_header(s, "03  ·  A TRIADE EXECUTORA",
                   "Tres operacoes de classe mundial. Uma unica visao.")

triade = [
    ("GSWP",
     "Incorporacao  ·  Estruturacao  ·  Governanca",
     "Sistema construtivo proprietario. S.A. com patrimonio de afetacao. "
     "Validada por fundos institucionais. Execucao end-to-end sob "
     "governanca auditada."),
    ("KONIGSBERGER VANNUCCHI",
     "Masterplan  ·  Arquitetura",
     "53 anos  ·  893 projetos  ·  R$ 18Bi de VGV  ·  107 arquitetos em 4 paises. "
     "Premio Master Imobiliario. MIPIM 2026 Cannes. 14a Bienal de Arquitetura."),
    ("HAMAM GLOBAL",
     "Hospitalidade do conceito a entrega",
     "26 anos  ·  +250 hoteis  ·  +70 marcas  ·  +43.000 chaves. "
     "Hubs em SP, Miami e Dubai. Eleita 6x melhor empresa de FF&E/OS&E "
     "da America Latina."),
]
cx, cy, cw = Cm(1.5), Cm(6.5), Cm(10.3)
for i, (nome, papel, desc) in enumerate(triade):
    x = cx + i * (cw + Cm(0.35))
    add_rect(s, x, cy, cw, Cm(11.3), fill=NAVY, radius=True)
    # Numero
    add_text(s, x + Cm(0.7), cy + Cm(0.5), Cm(2), Cm(2),
             f"0{i+1}", font=FONT_SERIF, size=46, color=TURQ,
             line_spacing=1.0)
    add_line(s, x + Cm(0.7), cy + Cm(3.1), x + Cm(2.8), cy + Cm(3.1),
             color=TURQ, weight=1.2)
    # Nome
    add_text(s, x + Cm(0.7), cy + Cm(3.4), cw - Cm(1.4), Cm(1.4),
             nome, font=FONT_SERIF, size=18, color=WHITE, bold=True,
             line_spacing=1.1)
    add_text(s, x + Cm(0.7), cy + Cm(5.0), cw - Cm(1.4), Cm(1),
             papel.upper(), font=FONT_SANS, size=8, color=TURQ,
             bold=True, letter_spacing=3)
    add_text(s, x + Cm(0.7), cy + Cm(6.2), cw - Cm(1.4), Cm(5),
             desc, font=FONT_SANS, size=10, color=CREAM, line_spacing=1.4)

add_footer(s, 4)


# ======================================================================
# BLOCO 02 - CONTEXTO MACRO: POR QUE AGORA
# ======================================================================

# --- S05 BRASIL NO RADAR GLOBAL -------------------------------------
s = new_slide(VOID)
add_section_header(s, "04  ·  CONTEXTO MACRO",
                   "O Brasil de volta ao radar do capital global.")

add_text(s, Cm(1.5), Cm(6.5), Cm(20), Cm(1.5),
         "BofA eleva projecao para o Ibovespa.",
         font=FONT_SERIF, size=24, color=WHITE, italic=True)
add_text(s, Cm(1.5), Cm(8.4), Cm(20), Cm(1),
         "Capital estrangeiro retorna a ativos reais brasileiros. "
         "Contexto favoravel a alternativos imobiliarios.",
         font=FONT_SANS, size=14, color=CREAM, line_spacing=1.4)

# Tres data points
data_macro = [
    ("Capital", "Estrangeiro", "Volta a ativos\nreais brasileiros"),
    ("BofA", "Bullish BR", "Projecao elevada\npara o Ibovespa"),
    ("HNWI", "Real Estate", "Demanda por\nalternativos premium"),
]
for i, (a, b, c) in enumerate(data_macro):
    x = Cm(1.5 + i * 7)
    y = Cm(11)
    add_rect(s, x, y, Cm(6.5), Cm(5.5), fill=NAVY, radius=True)
    add_rect(s, x, y, Cm(0.12), Cm(5.5), fill=TURQ)
    add_text(s, x + Cm(0.6), y + Cm(0.5), Cm(5.8), Cm(0.7),
             a.upper(), font=FONT_SANS, size=9, color=TURQ, bold=True,
             letter_spacing=4)
    add_text(s, x + Cm(0.6), y + Cm(1.2), Cm(5.8), Cm(1.5),
             b, font=FONT_SERIF, size=22, color=WHITE, bold=True)
    add_text(s, x + Cm(0.6), y + Cm(3.2), Cm(5.8), Cm(2.2),
             c, font=FONT_SANS, size=10, color=CREAM, line_spacing=1.4)

add_text(s, Cm(1.5), SLIDE_H - Cm(1.6), Cm(20), Cm(0.4),
         "Fonte: InfoMoney  ·  Bloomberg Linea",
         font=FONT_SANS, size=8, color=GREY_50, italic=True)
add_footer(s, 5)


# --- S06 NORDESTE EM ACELERACAO -------------------------------------
s = new_slide(VOID)
add_section_header(s, "05  ·  NORDESTE",
                   "Aceleracao estrutural sem precedentes.")

bullets_ne = [
    ("+20,85%", "Salvador", "Maior valorizacao entre capitais  ·  FipeZAP Mai/2025"),
    ("3 / 5", "Ranking Capitais", "Tres das cinco capitais mais valorizadas do pais sao do NE"),
    ("+85%", "Luxo & Superluxo NE", "Crescimento em apenas 2 anos  ·  Fonte: CBIC"),
    ("+80,3%", "Sergipe VGV", "Maior crescimento de VGV do Nordeste"),
    ("+20,5%", "Lancamentos 2024", "Crescimento regional ano contra ano"),
]
y = Cm(6.5)
for i, (val, label, desc) in enumerate(bullets_ne):
    yy = y + i * Cm(2.05)
    add_rect(s, Cm(1.5), yy, Cm(30.8), Cm(1.85), fill=NAVY, radius=True)
    add_text(s, Cm(2.0), yy + Cm(0.35), Cm(6), Cm(1.2),
             val, font=FONT_SERIF, size=26, color=TURQ, bold=True)
    add_text(s, Cm(8.5), yy + Cm(0.4), Cm(8), Cm(1.1),
             label.upper(), font=FONT_SANS, size=11, color=WHITE,
             bold=True, letter_spacing=3)
    add_text(s, Cm(16.5), yy + Cm(0.55), Cm(15.8), Cm(1.2),
             desc, font=FONT_SANS, size=10, color=CREAM)

add_footer(s, 6)


# --- S07 PARADOXO SERGIPANO -----------------------------------------
s = new_slide(VOID)
add_section_header(s, "06  ·  O PARADOXO SERGIPANO",
                   "A janela que se abre agora.")

# Quote
add_text(s, Cm(1.5), Cm(6.5), Cm(30), Cm(2),
         "\"Enquanto o Nordeste explodia de valor, Sergipe estava travada.\nEssa janela se abre agora.\"",
         font=FONT_SERIF, size=20, color=WHITE, italic=True, line_spacing=1.3)

# Tabela benchmark
add_text(s, Cm(1.5), Cm(9.5), Cm(20), Cm(0.7),
         "BENCHMARK  ·  SUPERLUXO FRENTE MAR  ·  NOVOS LANCAMENTOS",
         font=FONT_SANS, size=9, color=TURQ, bold=True, letter_spacing=4)

bench = [
    ("Salvador (BA)", "R$ 18.000 - 22.000 / m2", "R$ 28.000 - 35.000 / m2"),
    ("Maceio (AL)", "R$ 18.000 - 22.000 / m2", "R$ 28.000 - 35.000 / m2"),
    ("Praia do Forte (BA)", "R$ 20.000 - 25.000 / m2", "R$ 30.000 - 40.000 / m2"),
    ("Jurere Internacional (SC)", "R$ 25.000 - 30.000 / m2", "R$ 35.000+ / m2"),
    ("Aracaju beira-mar (novos)", "R$ 14.000 - 16.000 / m2", "—"),
]
ty = Cm(10.5)
# Header row
add_rect(s, Cm(1.5), ty, Cm(30.8), Cm(0.85), fill=NAVY_LIGHT)
add_text(s, Cm(2), ty + Cm(0.2), Cm(11), Cm(0.6), "MERCADO",
         font=FONT_SANS, size=9, color=TURQ, bold=True, letter_spacing=3)
add_text(s, Cm(13), ty + Cm(0.2), Cm(9), Cm(0.6), "ALTO PADRAO NOVOS",
         font=FONT_SANS, size=9, color=TURQ, bold=True, letter_spacing=3)
add_text(s, Cm(22), ty + Cm(0.2), Cm(10), Cm(0.6), "ALTO LUXO FRENTE MAR",
         font=FONT_SANS, size=9, color=TURQ, bold=True, letter_spacing=3)

for i, (mkt, hp, lux) in enumerate(bench):
    ry = ty + Cm(0.85) + i * Cm(0.85)
    if i % 2 == 0:
        add_rect(s, Cm(1.5), ry, Cm(30.8), Cm(0.85), fill=NAVY_LIGHT)
    add_text(s, Cm(2), ry + Cm(0.2), Cm(11), Cm(0.6), mkt,
             font=FONT_SERIF, size=12, color=WHITE)
    add_text(s, Cm(13), ry + Cm(0.22), Cm(9), Cm(0.6), hp,
             font=FONT_SANS, size=11, color=CREAM)
    add_text(s, Cm(22), ry + Cm(0.22), Cm(10), Cm(0.6), lux,
             font=FONT_SANS, size=11, color=CREAM)

# Callout
add_rect(s, Cm(1.5), Cm(15.7), Cm(30.8), Cm(2.2), fill=NAVY, radius=True)
add_rect(s, Cm(1.5), Cm(15.7), Cm(0.12), Cm(2.2), fill=GOLD)
add_text(s, Cm(2.2), Cm(16.0), Cm(29), Cm(1.7),
         "\"Estancia nao foi esquecida por falta de valor. Foi travada por "
         "fatores regulatorios e judiciais.\nO estado mais seguro do Nordeste — "
         "os holofotes voltam agora.\"",
         font=FONT_SERIF, size=13, color=WHITE, italic=True, line_spacing=1.4)
add_footer(s, 7)


# --- S08 OS 3 CATALISADORES -----------------------------------------
s = new_slide(VOID)
add_section_header(s, "07  ·  OS TRES CATALISADORES",
                   "O que destravou Estancia depois de uma decada.")

cats = [
    ("01", "JUDICIAL",
     "Acordo homologado pela\nJustica Federal",
     "Impasse de 10 anos encerrado.\nMPF + Gov. SE + Ibama + Adema + Prefeitura."),
    ("02", "REGULATORIO",
     "Novo Plano Diretor\naprovado",
     "De zona rural (lote minimo 20.000m2)\npara desenvolvimento urbano com\ngabarito de ate 7 pavimentos."),
    ("03", "MERCADO",
     "Aracaju beira-mar\na R$ 14k-16k/m2",
     "Correcao chegou a capital.\nO litoral sul ainda nao foi precificado."),
]
cy = Cm(6.5)
cw = Cm(10.3)
for i, (n, head, sub, body) in enumerate(cats):
    x = Cm(1.5) + i * (cw + Cm(0.35))
    add_rect(s, x, cy, cw, Cm(11.3), fill=NAVY, radius=True)
    add_text(s, x + Cm(0.7), cy + Cm(0.5), Cm(3), Cm(2.5),
             n, font=FONT_SERIF, size=60, color=TURQ, line_spacing=1.0)
    add_line(s, x + Cm(0.7), cy + Cm(3.5), x + cw - Cm(0.7), cy + Cm(3.5),
             color=TURQ, weight=1.0)
    add_text(s, x + Cm(0.7), cy + Cm(3.8), cw - Cm(1.4), Cm(0.8),
             head, font=FONT_SANS, size=10, color=TURQ, bold=True,
             letter_spacing=4)
    add_text(s, x + Cm(0.7), cy + Cm(4.8), cw - Cm(1.4), Cm(2.5),
             sub, font=FONT_SERIF, size=18, color=WHITE,
             bold=True, line_spacing=1.2)
    add_text(s, x + Cm(0.7), cy + Cm(7.5), cw - Cm(1.4), Cm(3.5),
             body, font=FONT_SANS, size=10, color=CREAM, line_spacing=1.5)
add_footer(s, 8)


# --- S09 PROJECAO DE DEMANDA / TAM ----------------------------------
s = new_slide(VOID)
add_section_header(s, "08  ·  PROJECAO DE DEMANDA / TAM",
                   "Quatro quadrantes de demanda convergem.")

quads = [
    ("LOCAL", "Raio 150 km",
     "~115.000 potenciais compradores  ·  primeira e segunda residencia"),
    ("REGIONAL", "Nordeste",
     "Salvador  ·  Recife  ·  Maceio  ·  destino premium fins de semana"),
    ("NACIONAL", "SP  ·  MG  ·  DF",
     "HNWIs  ·  second homes de inverno  ·  investidores de rentabilidade"),
    ("INTERNACIONAL", "Europa & Americas",
     "Turismo de luxo global  ·  nomades digitais premium"),
]
for i, (titulo, escopo, desc) in enumerate(quads):
    col = i % 2
    row = i // 2
    x = Cm(1.5) + col * Cm(15.7)
    y = Cm(6.5) + row * Cm(4.2)
    add_rect(s, x, y, Cm(15.4), Cm(3.9), fill=NAVY, radius=True)
    add_rect(s, x, y, Cm(0.12), Cm(3.9), fill=TURQ)
    add_text(s, x + Cm(0.7), y + Cm(0.4), Cm(14), Cm(0.7),
             titulo, font=FONT_SANS, size=10, color=TURQ, bold=True,
             letter_spacing=4)
    add_text(s, x + Cm(0.7), y + Cm(1.2), Cm(14), Cm(1.4),
             escopo, font=FONT_SERIF, size=22, color=WHITE, bold=True)
    add_text(s, x + Cm(0.7), y + Cm(2.7), Cm(14), Cm(1.2),
             desc, font=FONT_SANS, size=11, color=CREAM, line_spacing=1.3)

# Anchor
add_text(s, Cm(1.5), Cm(15.4), Cm(30.8), Cm(2),
         "\"Ecossistema continuo — atrai o comprador e o investidor — "
         "gerando ciclo virtuoso de demanda.\"",
         font=FONT_SERIF, size=15, color=GOLD, italic=True,
         align=PP_ALIGN.CENTER, line_spacing=1.3)
add_footer(s, 9)


# --- S10 BENCHMARKING FIRST MOVER -----------------------------------
s = new_slide(VOID)
add_section_header(s, "09  ·  BENCHMARKING",
                   "O efeito first mover comprovado.")

cols_bench = [
    ("JURERE INT.", "Consolidado", "4x - 5x / 10 anos",
     "Master plan + infraestrutura"),
    ("TRANCOSO", "Consolidado Premium", "3x - 4x / 15 anos",
     "Charme organico"),
    ("JERICOACOARA", "Emergente", "Acelerado",
     "Aeroporto catalisador"),
    ("JUBIABA", "A Nova Fronteira", "3x - 6x / 4-5 anos",
     "KV + Hamam + Smart City"),
]
cy = Cm(6.5)
cw = Cm(7.7)
for i, (nome, status, valoriz, perfil) in enumerate(cols_bench):
    x = Cm(1.5) + i * (cw + Cm(0.2))
    is_jub = (i == 3)
    fill = NAVY if not is_jub else TURQ
    add_rect(s, x, cy, cw, Cm(10.5), fill=fill, radius=True)
    txt_main = WHITE if not is_jub else VOID
    txt_sec = TURQ if not is_jub else VOID
    txt_body = CREAM if not is_jub else VOID
    add_text(s, x + Cm(0.5), cy + Cm(0.5), Cm(7), Cm(1),
             nome, font=FONT_SANS, size=10, color=txt_sec,
             bold=True, letter_spacing=3)
    add_line(s, x + Cm(0.5), cy + Cm(1.4), x + Cm(2.5), cy + Cm(1.4),
             color=txt_sec, weight=1.0)
    add_text(s, x + Cm(0.5), cy + Cm(1.7), Cm(7), Cm(2.5),
             status, font=FONT_SERIF, size=18, color=txt_main, bold=True,
             line_spacing=1.15)
    add_text(s, x + Cm(0.5), cy + Cm(4.7), Cm(7), Cm(0.7),
             "VALORIZACAO", font=FONT_SANS, size=8, color=txt_sec,
             bold=True, letter_spacing=3)
    add_text(s, x + Cm(0.5), cy + Cm(5.4), Cm(7), Cm(1.5),
             valoriz, font=FONT_SERIF, size=15, color=txt_main,
             line_spacing=1.2)
    add_text(s, x + Cm(0.5), cy + Cm(7.3), Cm(7), Cm(0.7),
             "PERFIL", font=FONT_SANS, size=8, color=txt_sec,
             bold=True, letter_spacing=3)
    add_text(s, x + Cm(0.5), cy + Cm(8.0), Cm(7), Cm(2.3),
             perfil, font=FONT_SANS, size=11, color=txt_body,
             line_spacing=1.4)

add_text(s, Cm(1.5), Cm(17.4), Cm(30.8), Cm(0.8),
         "\"O investidor da Fase 1 entra no momento de maxima assimetria "
         "— captura a curva que Jurere e Trancoso ja esgotaram.\"",
         font=FONT_SERIF, size=12, color=GOLD, italic=True,
         align=PP_ALIGN.CENTER)
add_footer(s, 10)


# --- S11 CONVERGENCIA: GRAFICO --------------------------------------
s = new_slide(VOID)
add_section_header(s, "10  ·  A CONVERGENCIA",
                   "O custo de conversao versus o mercado precificado.")

bars = [
    ("Jurere / Trancoso", 35000, GREY_50),
    ("Salvador / Maceio superluxo", 31500, GREY_50),
    ("Aracaju beira-mar (novos)", 15000, CREAM),
    ("JUBIABA  ·  custo de conversao", 4750, TURQ),
]
gx = Cm(1.5)
gy = Cm(7.0)
gw_total = Cm(28)
max_v = 35000
bar_h = Cm(1.4)
gap = Cm(0.55)

for i, (label, val, color) in enumerate(bars):
    yy = gy + i * (bar_h + gap)
    # label
    add_text(s, gx, yy + Cm(0.25), Cm(10), Cm(1),
             label, font=FONT_SANS, size=11, color=WHITE, bold=(i == 3))
    # bar bg
    bx = gx + Cm(10.3)
    bw = Cm(18)
    add_rect(s, bx, yy + Cm(0.45), bw, Cm(0.6), fill=NAVY)
    # bar fill
    fill_w = Emu(int(bw * val / max_v))
    add_rect(s, bx, yy + Cm(0.45), fill_w, Cm(0.6), fill=color)
    # value
    add_text(s, bx + bw + Cm(0.2), yy + Cm(0.25), Cm(5), Cm(1),
             f"R$ {val:,.0f}/m2".replace(",", "."),
             font=FONT_SERIF, size=13, color=color, bold=(i == 3))

add_rect(s, Cm(1.5), Cm(15.5), Cm(30.8), Cm(2.0), fill=NAVY, radius=True)
add_rect(s, Cm(1.5), Cm(15.5), Cm(0.12), Cm(2.0), fill=TURQ)
add_text(s, Cm(2.2), Cm(15.8), Cm(29), Cm(1.5),
         "Custo de conversao da cota ≠ preco de venda do imovel.\n"
         "Projecao de lancamento Fase 3:  R$ 14.000 - 16.000 / m2.",
         font=FONT_SERIF, size=12, color=WHITE, italic=True, line_spacing=1.4)
add_footer(s, 11)


# --- S12 A JANELA: AS 4 FASES ---------------------------------------
s = new_slide(VOID)
add_section_header(s, "11  ·  A JANELA QUE SE FECHA",
                   "As quatro fases do ciclo de Jubiaba.")

fases = [
    ("0", "CONCEPCAO", "[OK] Concluida",
     "GSWP assumiu 100% do risco inicial.", GREY_50),
    ("1", "DESENVOLVIMENTO", "Voce esta aqui",
     "Masterplan KV em execucao  ·  Hamam estruturada\n"
     "Maxima janela de entrada como socio fundador", TURQ),
    ("2", "PRE-LANCAMENTO", "Proxima",
     "Ativo registrado  ·  F&F e VIP\nAgio +30% a +40% vs. Fase 1", CREAM),
    ("3", "GO TO MARKET", "Futuro",
     "Lancamento publico\nR$ 14.000 - 16.000 / m2", GREY_50),
]
fy = Cm(7.2)
fw = Cm(7.4)
for i, (n, nome, status, desc, col) in enumerate(fases):
    x = Cm(1.5) + i * (fw + Cm(0.35))
    is_active = (i == 1)
    fill = TURQ if is_active else NAVY
    add_rect(s, x, fy, fw, Cm(8.5), fill=fill, radius=True)
    txt_main = VOID if is_active else WHITE
    txt_sec = VOID if is_active else TURQ
    txt_body = VOID if is_active else CREAM
    add_text(s, x + Cm(0.5), fy + Cm(0.4), Cm(2), Cm(2),
             n, font=FONT_SERIF, size=70, color=txt_main, line_spacing=1.0)
    add_text(s, x + Cm(0.5), fy + Cm(3.4), fw - Cm(1), Cm(0.7),
             nome, font=FONT_SANS, size=9, color=txt_sec, bold=True,
             letter_spacing=3)
    add_text(s, x + Cm(0.5), fy + Cm(4.2), fw - Cm(1), Cm(0.9),
             status, font=FONT_SERIF, size=14, color=txt_main, italic=True,
             bold=True)
    add_text(s, x + Cm(0.5), fy + Cm(5.4), fw - Cm(1), Cm(3),
             desc, font=FONT_SANS, size=10, color=txt_body, line_spacing=1.4)

add_text(s, Cm(1.5), Cm(16.6), Cm(30.8), Cm(0.8),
         "\"Quem entra na Fase 1 torna-se socio do projeto que a Fase 3 vai lancar.\"",
         font=FONT_SERIF, size=14, color=GOLD, italic=True,
         align=PP_ALIGN.CENTER)
add_footer(s, 12)


# ======================================================================
# BLOCO 03 - O ATIVO
# ======================================================================

# --- S13 A TERRA: NUMEROS -------------------------------------------
s = new_slide(VOID)
add_section_header(s, "12  ·  A TERRA",
                   "Os numeros de uma reserva irrepetivel.")

nums_terra = [
    ("129 ha", "Fase 1  ·  potencial +70 ha"),
    ("1.070 m", "Faixa de areia exclusiva pe na areia"),
    ("2.500 m", "Profundidade ate a Linha Verde (SE-100)"),
    ("15 min", "Centro industrial -> litoral (hoje 45 min)"),
    ("Mar + Rio + Lagoa", "Unico no litoral nordestino"),
    ("28 - 35%", "APP integrada como Parque Natural"),
]
gx = Cm(1.5)
gy = Cm(6.5)
gw = Cm(10.05)
gh = Cm(3.4)
for i, (num, sub) in enumerate(nums_terra):
    col = i % 3
    row = i // 3
    x = gx + col * (gw + Cm(0.4))
    y = gy + row * (gh + Cm(0.4))
    add_rect(s, x, y, gw, gh, fill=NAVY, radius=True)
    add_rect(s, x, y, Cm(0.12), gh, fill=TURQ)
    add_text(s, x + Cm(0.7), y + Cm(0.5), gw - Cm(0.9), Cm(1.5),
             num, font=FONT_SERIF, size=26, color=WHITE, bold=True)
    add_text(s, x + Cm(0.7), y + Cm(2.0), gw - Cm(0.9), Cm(1.3),
             sub, font=FONT_SANS, size=10, color=CREAM, line_spacing=1.4)

# Linha de identidade
add_text(s, Cm(1.5), Cm(14.3), Cm(30.8), Cm(2.2),
         "Praia do Saco (10km, top 100 mundial)  ·  Mangue Seco (12km, Tieta)\n"
         "Lagoa dos Tambaquis dentro da area  ·  Producao local: mangaba, caju, coco, acerola",
         font=FONT_SERIF, size=12, color=GOLD, italic=True,
         align=PP_ALIGN.CENTER, line_spacing=1.5)
add_footer(s, 13)


# --- S14 GALERIA I (full bleed) -------------------------------------
s = new_slide(VOID)
# Grade 2x3 de imagens
photos = [p for p in gallery_area if p][:6]
if not photos:
    photos = [img_aerea, img_aerea2, img_aerea3] * 2
gx = Cm(0)
gy = Cm(0)
cell_w = SLIDE_W / 3
cell_h = SLIDE_H / 2
for i, p in enumerate(photos[:6]):
    if not p:
        continue
    col = i % 3
    row = i // 3
    try:
        slide_pic = s.shapes.add_picture(p, col * cell_w, row * cell_h,
                                         width=cell_w, height=cell_h)
    except Exception:
        pass

# Overlay sutil + label
add_overlay(s, opacity=0.2, color=VOID)
add_text(s, Cm(1.5), Cm(0.8), Cm(20), Cm(0.7),
         "13  ·  GALERIA  ·  AREA",
         font=FONT_SANS, size=9, color=TURQ, bold=True, letter_spacing=4)
add_text(s, Cm(1.5), Cm(1.5), Cm(28), Cm(2),
         "1.070m de litoral virgem.",
         font=FONT_SERIF, size=36, color=WHITE)
add_footer(s, 14)


# --- S15 GALERIA II - NATUREZA --------------------------------------
s = new_slide(VOID)
photos2 = [p for p in gallery_regiao if p][:6]
for i, p in enumerate(photos2[:6]):
    if not p:
        continue
    col = i % 3
    row = i // 3
    try:
        s.shapes.add_picture(p, col * cell_w, row * cell_h,
                             width=cell_w, height=cell_h)
    except Exception:
        pass

add_overlay(s, opacity=0.2, color=VOID)
add_text(s, Cm(1.5), Cm(0.8), Cm(20), Cm(0.7),
         "14  ·  GALERIA  ·  NATUREZA",
         font=FONT_SANS, size=9, color=TURQ, bold=True, letter_spacing=4)
add_text(s, Cm(1.5), Cm(1.5), Cm(28), Cm(2),
         "Praia. Dunas. Lagoa. Mangue.",
         font=FONT_SERIF, size=36, color=WHITE)
add_footer(s, 15)


# --- S16 LOCALIZACAO ESTRATEGICA ------------------------------------
s = new_slide(VOID)
add_section_header(s, "15  ·  LOCALIZACAO ESTRATEGICA",
                   "A geografia do timing.")

# Esquerda - imagem aerea
if img_aero:
    try:
        s.shapes.add_picture(img_aero, Cm(1.5), Cm(6.5),
                             width=Cm(15), height=Cm(10.5))
    except Exception:
        pass

# Direita - dados
loc_x = Cm(17.5)
loc_y = Cm(6.5)
loc_data = [
    ("50 min", "Aeroporto de Aracaju  ·  voos diretos SP, BSB, RJ"),
    ("45 min", "Cidade de Aracaju"),
    ("Linha Verde", "Acesso direto SE-100"),
    ("10 km", "Praia do Saco  ·  top 100 mundial"),
    ("12 km", "Mangue Seco  ·  novela Tieta"),
    ("10 km", "Divisa com a Bahia"),
    ("Estancia", "Polo industrial  ·  primeira residencia"),
]
add_rect(s, loc_x, loc_y, Cm(14.8), Cm(10.5), fill=NAVY, radius=True)
add_text(s, loc_x + Cm(0.6), loc_y + Cm(0.4), Cm(13), Cm(0.7),
         "RAIO DE INFLUENCIA",
         font=FONT_SANS, size=10, color=TURQ, bold=True, letter_spacing=4)
for i, (a, b) in enumerate(loc_data):
    yy = loc_y + Cm(1.4) + i * Cm(1.25)
    add_text(s, loc_x + Cm(0.6), yy, Cm(4.5), Cm(0.8),
             a, font=FONT_SERIF, size=14, color=TURQ, bold=True)
    add_text(s, loc_x + Cm(5.5), yy + Cm(0.05), Cm(9), Cm(0.8),
             b, font=FONT_SANS, size=10, color=CREAM)

add_text(s, Cm(1.5), Cm(17.4), Cm(30.8), Cm(0.6),
         "Infraestrutura adjacente: Porto de Sergipe  ·  Iate Clube  ·  "
         "Heliponto  ·  Marina Porto N'Angola",
         font=FONT_SERIF, size=11, color=GOLD, italic=True,
         align=PP_ALIGN.CENTER)
add_footer(s, 16)


# --- S17 MASTERPLAN KV - ZONEAMENTO ---------------------------------
s = new_slide(VOID)
add_section_header(s, "16  ·  O MASTERPLAN  ·  KONIGSBERGER VANNUCCHI",
                   "Zoneamento oficial ZEEC / Estancia  ·  Lei 8.980/2022.")

zones = [
    ("A", "Parque Natural de Dunas & Restinga", "ZPA / ZMAR", "28 - 35%"),
    ("B", "Beach Club + Hotel / Condo-Hotel", "ZUR (fora de APP)", "6 - 9%"),
    ("C", "Condo Village  ·  Apartamentos", "ZUD / ZUC", "10 - 13%"),
    ("D", "Bairros de Lotes Residenciais", "ZUD / ZUC", "32 - 40%"),
    ("E", "Boulevard Comercial & Esportivo", "ZUD / ZUC", "3 - 5%"),
    ("F", "Utilidades / Infraestrutura", "ZUE / ZUD", "2 - 3%"),
]

ty = Cm(6.5)
# Cabecalho
add_rect(s, Cm(1.5), ty, Cm(30.8), Cm(0.85), fill=NAVY_LIGHT)
heads = [("ZONA", Cm(2)), ("USO", Cm(4)), ("CLASSIFICACAO", Cm(20)),
         ("% ABTA", Cm(27))]
for h, hx in heads:
    add_text(s, hx, ty + Cm(0.2), Cm(8), Cm(0.6), h,
             font=FONT_SANS, size=9, color=TURQ, bold=True, letter_spacing=3)

for i, (z, u, c, p) in enumerate(zones):
    ry = ty + Cm(0.85) + i * Cm(1.05)
    if i % 2 == 0:
        add_rect(s, Cm(1.5), ry, Cm(30.8), Cm(1.05), fill=NAVY_LIGHT)
    add_text(s, Cm(2), ry + Cm(0.3), Cm(2), Cm(0.6),
             z, font=FONT_SERIF, size=20, color=TURQ, bold=True)
    add_text(s, Cm(4), ry + Cm(0.32), Cm(15.5), Cm(0.6),
             u, font=FONT_SERIF, size=13, color=WHITE)
    add_text(s, Cm(20), ry + Cm(0.35), Cm(7), Cm(0.6),
             c, font=FONT_SANS, size=10, color=CREAM)
    add_text(s, Cm(27), ry + Cm(0.3), Cm(5), Cm(0.6),
             p, font=FONT_SERIF, size=14, color=TURQ, bold=True)

add_text(s, Cm(1.5), Cm(17.6), Cm(30.8), Cm(0.6),
         "Orla intocada — sem arruamento nem volumes permanentes na faixa "
         "de praia. Requisito legal e argumento central de marketing.",
         font=FONT_SERIF, size=11, color=GOLD, italic=True,
         align=PP_ALIGN.CENTER)
add_footer(s, 17)


# --- S18 SEASIDE + HINTERLAND ---------------------------------------
s = new_slide(VOID)
add_section_header(s, "17  ·  AS DUAS ETAPAS",
                   "Seaside e Hinterland  ·  da praia ao interior do terreno.")

etapas = [
    ("ETAPA 01", "JUBIABA SEASIDE",
     "Primeira frente de mar.\n\n"
     "•  Lotes premium\n"
     "•  Beach Villas\n"
     "•  Beach Club ancora\n"
     "•  Planejamento de baixissima densidade",
     TURQ),
    ("ETAPA 02", "JUBIABA HINTERLAND",
     "Interior do terreno.\n\n"
     "•  Hotel / resort frente mar\n"
     "•  Hotel boutique\n"
     "•  Hotel a beira da Linha Verde\n"
     "•  Wellness  ·  Condo-hotel",
     CREAM),
]
for i, (label, nome, desc, accent) in enumerate(etapas):
    x = Cm(1.5) + i * Cm(15.7)
    add_rect(s, x, Cm(6.5), Cm(15.4), Cm(8.7), fill=NAVY, radius=True)
    add_rect(s, x, Cm(6.5), Cm(0.12), Cm(8.7), fill=accent)
    add_text(s, x + Cm(0.7), Cm(7.0), Cm(14), Cm(0.7),
             label, font=FONT_SANS, size=10, color=accent, bold=True,
             letter_spacing=4)
    add_text(s, x + Cm(0.7), Cm(7.8), Cm(14), Cm(1.5),
             nome, font=FONT_SERIF, size=24, color=WHITE, bold=True)
    add_text(s, x + Cm(0.7), Cm(9.7), Cm(14), Cm(5.2),
             desc, font=FONT_SANS, size=12, color=CREAM, line_spacing=1.5)

# Villas Suspensas highlight
add_rect(s, Cm(1.5), Cm(15.7), Cm(30.8), Cm(2.0), fill=NAVY_LIGHT, radius=True)
add_rect(s, Cm(1.5), Cm(15.7), Cm(0.12), Cm(2.0), fill=GOLD)
add_text(s, Cm(2.2), Cm(15.95), Cm(8), Cm(0.6),
         "VILLAS SUSPENSAS",
         font=FONT_SANS, size=10, color=GOLD, bold=True, letter_spacing=4)
add_text(s, Cm(2.2), Cm(16.55), Cm(28.5), Cm(1.2),
         "Arquitetura elevada que toca o solo de forma pontual — preserva a "
         "vegetacao sob a estrutura, garante vistas panoramicas, com zero "
         "movimentacao de terra.",
         font=FONT_SERIF, size=11, color=CREAM, italic=True, line_spacing=1.4)
add_footer(s, 18)


# --- S19 PILARES DO COMPLEXO ----------------------------------------
s = new_slide(VOID)
add_section_header(s, "18  ·  OS PILARES DO COMPLEXO",
                   "Quatro frentes de geracao de valor.")

pilares = [
    ("RESIDENCIAL", "Beach Villas frente mar  +  Condo Village",
     "100 - 140 m2  ·  parte com lock-off / pool. Sistema GSWP — "
     "qualidade superior, ate 3x mais rapido que metodo convencional."),
    ("HOTELARIA", "Resort frente mar (60-80 chaves)  +  Boutique  +  Linha Verde",
     "Padrao Hamam Global. Operacao IAAS. FF&E e OS&E end-to-end. "
     "Pool de marcas: Marriott, Hilton, Accor, IHG, Oetker."),
    ("LIFESTYLE & WELLNESS",
     "Beach Club  ·  Spa  ·  Centro Esportivo  ·  Tenis  ·  Beach Tennis",
     "Alta gastronomia farm-to-table  ·  Jardins sensoriais  ·  "
     "Trilhas ecologicas  ·  Wellness Center."),
    ("SUSTENTABILIDADE",
     "APP como Parque Natural  ·  Dark-sky (tartarugas marinhas)",
     "Energia solar  ·  SUDS  ·  gestao de efluentes  ·  zero poluicao "
     "visual  ·  impacto local positivo."),
]
for i, (titulo, sub, desc) in enumerate(pilares):
    col = i % 2
    row = i // 2
    x = Cm(1.5) + col * Cm(15.7)
    y = Cm(6.5) + row * Cm(5.3)
    add_rect(s, x, y, Cm(15.4), Cm(5.0), fill=NAVY, radius=True)
    add_rect(s, x, y, Cm(0.12), Cm(5.0), fill=TURQ)
    add_text(s, x + Cm(0.6), y + Cm(0.4), Cm(14), Cm(0.7),
             titulo, font=FONT_SANS, size=10, color=TURQ, bold=True,
             letter_spacing=4)
    add_text(s, x + Cm(0.6), y + Cm(1.2), Cm(14), Cm(1.4),
             sub, font=FONT_SERIF, size=14, color=WHITE, bold=True,
             line_spacing=1.2)
    add_text(s, x + Cm(0.6), y + Cm(2.8), Cm(14), Cm(2.2),
             desc, font=FONT_SANS, size=10, color=CREAM, line_spacing=1.4)

add_footer(s, 19)


# --- S20 SMART CITY - INFRA INTELIGENTE ------------------------------
s = new_slide(VOID)
add_section_header(s, "19  ·  SMART CITY",
                   "Seis pilares integrados desde o zero.")

smart = [
    ("ENERGIA LIMPA",
     "Solar + eolica  ·  matriz renovavel"),
    ("INFRA ENTERRADA",
     "Cabeamento 100% subterraneo\nzero poluicao visual"),
    ("CONECTIVIDADE TOTAL",
     "Fibra optica  ·  automacao residencial\ntrabalho remoto premium"),
    ("MOBILIDADE INTELIGENTE",
     "EV chargers  ·  shuttle eletrico\nciclofaixas"),
    ("EFICIENCIA HIDRICA",
     "Captacao pluvial  ·  reuso\nirrigacao inteligente"),
    ("ARQUITETURA ORGANICA",
     "Estudos de insolacao e ventos\nconforto termico natural"),
]
gx = Cm(1.5)
gy = Cm(6.5)
gw = Cm(10.05)
gh = Cm(5.4)
for i, (titulo, desc) in enumerate(smart):
    col = i % 3
    row = i // 3
    x = gx + col * (gw + Cm(0.4))
    y = gy + row * (gh + Cm(0.4))
    add_rect(s, x, y, gw, gh, fill=NAVY, radius=True)
    add_text(s, x + Cm(0.7), y + Cm(0.5), Cm(2), Cm(2),
             f"0{i+1}", font=FONT_SERIF, size=46, color=TURQ,
             line_spacing=1.0)
    add_line(s, x + Cm(0.7), y + Cm(3.0), x + Cm(2.8), y + Cm(3.0),
             color=TURQ, weight=1.0)
    add_text(s, x + Cm(0.7), y + Cm(3.2), gw - Cm(1), Cm(0.8),
             titulo, font=FONT_SANS, size=10, color=WHITE,
             bold=True, letter_spacing=3)
    add_text(s, x + Cm(0.7), y + Cm(4.0), gw - Cm(1), Cm(1.4),
             desc, font=FONT_SANS, size=9, color=CREAM, line_spacing=1.35)
add_footer(s, 20)


# --- S21 COMPARATIVO DE ESCALA --------------------------------------
s = new_slide(VOID)
add_section_header(s, "20  ·  COMPARATIVO DE ESCALA",
                   "Onde Jubiaba se posiciona no mercado nacional.")

comp = [
    ("Jurere Internacional (SC)", "Superluxo frente mar", "R$ 28.000 - 35.000 / m2"),
    ("Praia do Forte (BA)", "Alto luxo frente mar", "R$ 30.000 - 40.000 / m2"),
    ("Destino Baixio (BA)", "Branded residences", "R$ 600K - R$ 2M+ (lote / villa)"),
    ("Alexandria Trancoso (BA)", "Alto padrao", "R$ 614K+ (lote)"),
    ("Noah Home & Spa (Caueira/SE)", "Concorrente regional", "R$ 300 - 500K (lotes)"),
    ("Belleville Litoral Sul (SE)", "Concorrente regional", "R$ 250K+ (lotes)"),
    ("JUBIABA  ·  FASE 3", "Superluxo pe na areia", "R$ 14.000 - 16.000 / m2"),
]
ty = Cm(6.5)
add_rect(s, Cm(1.5), ty, Cm(30.8), Cm(0.85), fill=NAVY_LIGHT)
hd = [("EMPREENDIMENTO", Cm(2)), ("TIPOLOGIA", Cm(15)),
      ("PRECO REFERENCIA", Cm(23))]
for h, hx in hd:
    add_text(s, hx, ty + Cm(0.2), Cm(13), Cm(0.6), h,
             font=FONT_SANS, size=9, color=TURQ, bold=True, letter_spacing=3)

for i, (e, t, p) in enumerate(comp):
    ry = ty + Cm(0.85) + i * Cm(1.0)
    is_jub = (i == len(comp) - 1)
    if is_jub:
        add_rect(s, Cm(1.5), ry, Cm(30.8), Cm(1.0), fill=TURQ)
        col1 = VOID
        col2 = VOID
        col3 = VOID
        f1 = FONT_SERIF
    else:
        if i % 2 == 0:
            add_rect(s, Cm(1.5), ry, Cm(30.8), Cm(1.0), fill=NAVY_LIGHT)
        col1 = WHITE
        col2 = CREAM
        col3 = CREAM
        f1 = FONT_SERIF
    add_text(s, Cm(2), ry + Cm(0.3), Cm(13), Cm(0.6),
             e, font=f1, size=12, color=col1, bold=is_jub)
    add_text(s, Cm(15), ry + Cm(0.32), Cm(8), Cm(0.6),
             t, font=FONT_SANS, size=10, color=col2)
    add_text(s, Cm(23), ry + Cm(0.3), Cm(9), Cm(0.6),
             p, font=FONT_SERIF, size=12, color=col3, bold=is_jub)
add_footer(s, 21)


# ======================================================================
# BLOCO 03B - OS PARCEIROS
# ======================================================================

# --- S22 KV ---------------------------------------------------------
s = new_slide(VOID)
add_section_header(s, "21  ·  KONIGSBERGER VANNUCCHI",
                   "A assinatura do projeto.")

# Esquerda - tese
tese_kv = (
    "Por que a KV para Jubiaba?\n\n"
    "Sao os unicos no Brasil com portfolio combinando masterplans de "
    "grande escala, residencial de altissimo padrao e expertise comprovada "
    "no Norte e Nordeste.\n\n"
    "A chancela KV elimina o risco de design e adiciona valor estetico "
    "inquestionavel — diferencial direto na precificacao."
)
add_text(s, Cm(1.5), Cm(6.5), Cm(15), Cm(11), tese_kv,
         font=FONT_SERIF, size=14, color=CREAM, line_spacing=1.5)

# Direita - credenciais
credx = Cm(17.5)
credy = Cm(6.5)
add_rect(s, credx, credy, Cm(14.8), Cm(11), fill=NAVY, radius=True)
add_text(s, credx + Cm(0.6), credy + Cm(0.4), Cm(13), Cm(0.7),
         "CREDENCIAIS", font=FONT_SANS, size=10, color=TURQ,
         bold=True, letter_spacing=4)

cred_kv = [
    ("1972", "Fundacao  ·  53 anos de operacao"),
    ("893", "Projetos  ·  516 construidos"),
    ("R$ 18Bi", "VGV desde 2022  ·  +7MM m2 na decada"),
    ("107", "Arquitetos em 4 paises  ·  1,2MM h em BIM"),
    ("30o", "Premio Master Imobiliario"),
    ("MIPIM", "Cannes 2026  ·  14a Bienal Arquitetura"),
]
for i, (a, b) in enumerate(cred_kv):
    yy = credy + Cm(1.4) + i * Cm(1.4)
    add_text(s, credx + Cm(0.6), yy, Cm(4.5), Cm(0.9),
             a, font=FONT_SERIF, size=18, color=TURQ, bold=True)
    add_text(s, credx + Cm(5.5), yy + Cm(0.18), Cm(9), Cm(0.9),
             b, font=FONT_SANS, size=10, color=CREAM)

add_text(s, Cm(1.5), Cm(17.4), Cm(30.8), Cm(0.6),
         "Projetos de referencia: Platina 220 SP  ·  Brascan Century Plaza  "
         "·  Sesc Paulista  ·  HY Pinheiros  ·  Insigna RJ",
         font=FONT_SERIF, size=10, color=GOLD, italic=True,
         align=PP_ALIGN.CENTER)
add_footer(s, 22)


# --- S23 HAMAM GLOBAL -----------------------------------------------
s = new_slide(VOID)
add_section_header(s, "22  ·  HAMAM GLOBAL",
                   "A garantia da hospitalidade.")

tese_hm = (
    "Por que a Hamam Global para Jubiaba?\n\n"
    "Sao os unicos no Brasil capazes de operar do conceito ao FF&E/OS&E com "
    "padrao internacional. O modelo IAAS elimina o risco de execucao "
    "hoteleira.\n\n"
    "Mais de 1.500 fabricantes parceiros = acabamento internacional a custo "
    "competitivo. Marcas top tier ja homologadas."
)
add_text(s, Cm(1.5), Cm(6.5), Cm(15), Cm(11), tese_hm,
         font=FONT_SERIF, size=14, color=CREAM, line_spacing=1.5)

credx = Cm(17.5)
credy = Cm(6.5)
add_rect(s, credx, credy, Cm(14.8), Cm(11), fill=NAVY, radius=True)
add_text(s, credx + Cm(0.6), credy + Cm(0.4), Cm(13), Cm(0.7),
         "CREDENCIAIS", font=FONT_SANS, size=10, color=TURQ,
         bold=True, letter_spacing=4)
cred_hm = [
    ("26", "Anos de operacao  ·  hubs SP, Miami, Dubai"),
    ("+250", "Hoteis entregues  ·  +70 marcas"),
    ("+43.000", "Chaves  ·  US$ 215MM em FF&E/OS&E"),
    ("+1.500", "Fabricantes parceiros homologados"),
    ("6x", "Melhor empresa FF&E/OS&E America Latina"),
    ("1875", "Reconhecida pela Hotel Management Magazine"),
]
for i, (a, b) in enumerate(cred_hm):
    yy = credy + Cm(1.4) + i * Cm(1.4)
    add_text(s, credx + Cm(0.6), yy, Cm(4.5), Cm(0.9),
             a, font=FONT_SERIF, size=18, color=TURQ, bold=True)
    add_text(s, credx + Cm(5.5), yy + Cm(0.18), Cm(9), Cm(0.9),
             b, font=FONT_SANS, size=10, color=CREAM)

add_text(s, Cm(1.5), Cm(17.4), Cm(30.8), Cm(0.6),
         "Marcas homologadas: Marriott  ·  Hilton  ·  Accor  ·  IHG  ·  "
         "Oetker Collection  ·  W Hotels  ·  Fasano",
         font=FONT_SERIF, size=10, color=GOLD, italic=True,
         align=PP_ALIGN.CENTER)
add_footer(s, 23)


# --- S24 HOSPITALIDADE INTEGRADA ------------------------------------
s = new_slide(VOID)
add_section_header(s, "23  ·  HOSPITALIDADE INTEGRADA",
                   "Tres camadas que transformam residencial em experiencia.")

camadas = [
    ("CAMADA 01", "SERVICOS ESSENCIAIS",
     "Limpeza  ·  manutencao  ·  jardinagem de precisao.\n"
     "Operacao end-to-end sob padrao internacional."),
    ("CAMADA 02", "GASTRONOMIA PREMIUM",
     "Restaurante signature  ·  room service 24/7  ·  "
     "menu assinado por chef de renome."),
    ("CAMADA 03", "O TOQUE JUBIABA",
     "Mimos diarios  ·  frutas e agua de coco na praia  ·  "
     "curadoria de experiencias locais autenticas."),
]
for i, (label, nome, desc) in enumerate(camadas):
    x = Cm(1.5) + i * Cm(10.5)
    add_rect(s, x, Cm(7.0), Cm(10.2), Cm(10.3), fill=NAVY, radius=True)
    # Big number
    add_text(s, x + Cm(0.6), Cm(7.5), Cm(3), Cm(2.5),
             f"0{i+1}", font=FONT_SERIF, size=72, color=TURQ, line_spacing=1.0)
    add_line(s, x + Cm(0.6), Cm(10.6), x + Cm(3.5), Cm(10.6),
             color=TURQ, weight=1.0)
    add_text(s, x + Cm(0.6), Cm(10.9), Cm(9), Cm(0.7),
             label, font=FONT_SANS, size=10, color=TURQ,
             bold=True, letter_spacing=4)
    add_text(s, x + Cm(0.6), Cm(11.7), Cm(9), Cm(2),
             nome, font=FONT_SERIF, size=18, color=WHITE, bold=True,
             line_spacing=1.15)
    add_text(s, x + Cm(0.6), Cm(14.2), Cm(9), Cm(3),
             desc, font=FONT_SANS, size=11, color=CREAM, line_spacing=1.45)

add_footer(s, 24)


# ======================================================================
# BLOCO 04 - MODELO DE NEGOCIO & FINANCEIRO
# ======================================================================

# --- S25 O QUE VOCE ESTA COMPRANDO ----------------------------------
s = new_slide(VOID)
add_section_header(s, "24  ·  O QUE VOCE ESTA COMPRANDO",
                   "Clareza conceitual fundamental.")

itens25 = [
    ("01", "NAO E A COMPRA DE UM IMOVEL",
     "E uma participacao na fase de desenvolvimento que se "
     "converte em m2 com desagio."),
    ("02", "VOCE ENTRA COMO SOCIO FUNDADOR",
     "Fase de maximo potencial — antes da precificacao publica. "
     "Captura integral da margem do incorporador."),
    ("03", "A CONVERSAO SE DA COM O RI",
     "Registro de Incorporacao formaliza a conversao da cota em "
     "metro quadrado escriturado."),
    ("04", "TRES OPCOES DE SAIDA",
     "Venda da cota  ·  recebimento do imovel  ·  permanencia no "
     "pool de locacao hoteleira gerido pela Hamam Global."),
]
for i, (n, t, d) in enumerate(itens25):
    col = i % 2
    row = i // 2
    x = Cm(1.5) + col * Cm(15.7)
    y = Cm(6.5) + row * Cm(5.4)
    add_rect(s, x, y, Cm(15.4), Cm(5.1), fill=NAVY, radius=True)
    add_text(s, x + Cm(0.6), y + Cm(0.4), Cm(2), Cm(2),
             n, font=FONT_SERIF, size=42, color=TURQ, line_spacing=1.0)
    add_text(s, x + Cm(3.5), y + Cm(0.7), Cm(11.5), Cm(1.4),
             t, font=FONT_SANS, size=11, color=WHITE, bold=True,
             letter_spacing=2)
    add_text(s, x + Cm(3.5), y + Cm(2.2), Cm(11.5), Cm(2.8),
             d, font=FONT_SERIF, size=12, color=CREAM, line_spacing=1.4)
add_footer(s, 25)


# --- S26 MODELO DE NEGOCIO: 5 ETAPAS --------------------------------
s = new_slide(VOID)
add_section_header(s, "25  ·  CICLO DE GERACAO DE VALOR",
                   "Cinco etapas. Uma S.A. blindada.")

etapas26 = [
    ("01", "PRE-LANCAMENTO",
     "Captacao seed capital  ·  masterplan KV  ·  aprovacoes ZEEC"),
    ("02", "DESENVOLVIMENTO",
     "Execucao de obras de infra  ·  sistema construtivo GSWP"),
    ("03", "PRE-VENDA",
     "Abertura de vendas com precificacao premium  ·  fluxo de caixa"),
    ("04", "ENTREGA",
     "Conclusao das obras  ·  conversao de cotas em m2 (3x - 6x)"),
    ("05", "POS-ENTREGA",
     "Operacao Hamam  ·  receita recorrente passiva para investidores"),
]
ey = Cm(7.5)
ew = Cm(6.05)
for i, (n, t, d) in enumerate(etapas26):
    x = Cm(1.5) + i * (ew + Cm(0.2))
    add_rect(s, x, ey, ew, Cm(8.5), fill=NAVY, radius=True)
    # Number circle
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Cm(0.5), ey + Cm(0.5),
                              Cm(2.0), Cm(2.0))
    circ.fill.solid()
    circ.fill.fore_color.rgb = TURQ
    circ.line.fill.background()
    circ.shadow.inherit = False
    add_text(s, x + Cm(0.5), ey + Cm(0.65), Cm(2), Cm(1.7),
             n, font=FONT_SERIF, size=22, color=VOID, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, x + Cm(0.4), ey + Cm(2.9), ew - Cm(0.8), Cm(1.4),
             t, font=FONT_SANS, size=10, color=WHITE,
             bold=True, letter_spacing=2)
    add_text(s, x + Cm(0.4), ey + Cm(4.5), ew - Cm(0.8), Cm(3.8),
             d, font=FONT_SANS, size=9, color=CREAM, line_spacing=1.4)
    if i < 4:
        # arrow
        add_text(s, x + ew - Cm(0.05), ey + Cm(1.1), Cm(0.5), Cm(1),
                 ">", font=FONT_SERIF, size=22, color=TURQ, bold=True)
add_footer(s, 26)


# --- S27 ESTRUTURA SOCIETARIA ---------------------------------------
s = new_slide(VOID)
add_section_header(s, "26  ·  ESTRUTURA SOCIETARIA & GARANTIA REAL",
                   "A blindagem mais solida do mercado.")

itens27 = [
    ("S.A. com Patrimonio de Afetacao",
     "100% dedicado ao projeto  ·  blindado contra terceiros"),
    ("Governanca Auditada Externamente",
     "Relatorios trimestrais  ·  due diligence permanente"),
    ("Garantia Real",
     "Investimento garantido pela matricula do terreno ate o RI"),
    ("Seguro Garantia de Obras",
     "Garante que a obra comeca e termina  ·  seguradora independente"),
    ("Validacao Institucional",
     "Quando ha fundo, ha garantia de execucao  ·  selo de seriedade"),
    ("Registros Formais em Cartorio",
     "Cada etapa registrada  ·  zero ambiguidade juridica"),
]
gx = Cm(1.5)
gy = Cm(6.5)
gw = Cm(15.4)
gh = Cm(1.85)
for i, (t, d) in enumerate(itens27):
    col = i % 2
    row = i // 2
    x = gx + col * (gw + Cm(0.3))
    y = gy + row * (gh + Cm(0.3))
    add_rect(s, x, y, gw, gh, fill=NAVY, radius=True)
    add_rect(s, x, y, Cm(0.12), gh, fill=TURQ)
    add_text(s, x + Cm(0.6), y + Cm(0.3), gw - Cm(0.8), Cm(0.8),
             t, font=FONT_SERIF, size=14, color=WHITE, bold=True)
    add_text(s, x + Cm(0.6), y + Cm(1.05), gw - Cm(0.8), Cm(0.8),
             d, font=FONT_SANS, size=10, color=CREAM)
add_footer(s, 27)


# --- S28 SISTEMA CONSTRUTIVO GSWP -----------------------------------
s = new_slide(VOID)
add_section_header(s, "27  ·  O SISTEMA CONSTRUTIVO GSWP",
                   "Parede concretada in loco  ·  diferencial direto na TIR.")

# Big stat
add_rect(s, Cm(1.5), Cm(6.5), Cm(14), Cm(11), fill=NAVY, radius=True)
add_rect(s, Cm(1.5), Cm(6.5), Cm(0.12), Cm(11), fill=TURQ)
add_text(s, Cm(2.2), Cm(7.0), Cm(13), Cm(1),
         "VELOCIDADE", font=FONT_SANS, size=10, color=TURQ,
         bold=True, letter_spacing=4)
add_text(s, Cm(2.2), Cm(7.8), Cm(13), Cm(5),
         "ATE 3X\nMAIS RAPIDO",
         font=FONT_SERIF, size=58, color=WHITE, bold=True, line_spacing=1.0)
add_text(s, Cm(2.2), Cm(13.5), Cm(13), Cm(3.5),
         "Que o metodo construtivo convencional. Estrutura ja nasce pronta "
         "para receber acabamento premium.",
         font=FONT_SERIF, size=12, color=CREAM, italic=True,
         line_spacing=1.4)

# Bullets right
bx = Cm(16.0)
by = Cm(6.5)
bw_card = Cm(16.3)
add_rect(s, bx, by, bw_card, Cm(11), fill=NAVY, radius=True)
add_text(s, bx + Cm(0.6), by + Cm(0.4), Cm(15), Cm(0.7),
         "VANTAGENS DIRETAS",
         font=FONT_SANS, size=10, color=TURQ, bold=True, letter_spacing=4)
vant = [
    ("Qualidade Superior", "Curadoria de alto padrao end-to-end"),
    ("Estrutura Otimizada", "Pronta para acabamento premium"),
    ("Custos Controlados", "Previsibilidade financeira absoluta"),
    ("Sustentabilidade", "Reducao de entulho e desperdicio"),
    ("Diferencial na TIR", "Impacto direto no retorno do investidor"),
]
for i, (t, d) in enumerate(vant):
    yy = by + Cm(1.4) + i * Cm(1.85)
    add_text(s, bx + Cm(0.6), yy, Cm(6), Cm(0.7),
             t, font=FONT_SERIF, size=14, color=WHITE, bold=True)
    add_text(s, bx + Cm(0.6), yy + Cm(0.85), Cm(15), Cm(0.7),
             d, font=FONT_SANS, size=10, color=CREAM)
add_footer(s, 28)


# --- S29 CICLO COMPLETO DE INVESTIMENTO -----------------------------
s = new_slide(VOID)
add_section_header(s, "28  ·  CICLO COMPLETO DE INVESTIMENTO",
                   "Garantia real ativa em todas as etapas.")

steps = [
    ("APORTE", "Capital aportado\nna S.A."),
    ("DESENVOLVIMENTO", "Obras de\ninfraestrutura"),
    ("RI", "Registro de\nIncorporacao"),
    ("CONVERSAO", "Cota convertida\nem m2"),
    ("SAIDA", "Venda  ·  uso  ·\ncondo-hotel"),
]
sy = Cm(8.5)
sw = Cm(6.05)
for i, (t, d) in enumerate(steps):
    x = Cm(1.5) + i * (sw + Cm(0.2))
    add_rect(s, x, sy, sw, Cm(5), fill=NAVY, radius=True)
    add_text(s, x + Cm(0.4), sy + Cm(0.4), sw - Cm(0.8), Cm(1),
             f"M{['0','1','2','3','4'][i]}", font=FONT_SANS, size=10,
             color=TURQ, bold=True, letter_spacing=4)
    add_text(s, x + Cm(0.4), sy + Cm(1.2), sw - Cm(0.8), Cm(1),
             t, font=FONT_SERIF, size=14, color=WHITE, bold=True)
    add_text(s, x + Cm(0.4), sy + Cm(2.5), sw - Cm(0.8), Cm(2),
             d, font=FONT_SANS, size=10, color=CREAM, line_spacing=1.4)
    if i < 4:
        # connector arrow
        add_text(s, x + sw - Cm(0.1), sy + Cm(1.7), Cm(0.5), Cm(1),
                 ">", font=FONT_SERIF, size=22, color=TURQ, bold=True)

# Garantia faixa
add_rect(s, Cm(1.5), Cm(14.5), Cm(30.8), Cm(2.3), fill=TURQ, radius=True)
add_text(s, Cm(2.2), Cm(14.85), Cm(29), Cm(1),
         "GARANTIA REAL", font=FONT_SANS, size=11, color=VOID,
         bold=True, letter_spacing=4)
add_text(s, Cm(2.2), Cm(15.6), Cm(29), Cm(1.2),
         "Imobiliaria ativa em todas as fases — ate a conversao em ativo registrado.",
         font=FONT_SERIF, size=14, color=VOID, bold=True, italic=True)
add_footer(s, 29)


# --- S30 ESTRATEGIA DE COMERCIALIZACAO ------------------------------
s = new_slide(VOID)
add_section_header(s, "29  ·  ESTRATEGIA DE COMERCIALIZACAO",
                   "Quatro fases com curva de precificacao crescente.")

fases30 = [
    ("FASE 01", "DESENVOLVIMENTO", "Investidores Semente",
     "R$ 4.000 - 5.497 / m2", "Custo base", TURQ),
    ("FASE 02", "PRE-LANCAMENTO", "F&F + Clientes VIP",
     "Agio +30 a +40%", "vs. Fase 1", CREAM),
    ("FASE 03", "LANCAMENTO", "Nacional / Internacional",
     "R$ 14.000 - 16.000 / m2", "Mercado publico", GREY_50),
    ("FASE 04", "OPERACAO", "Turistas High-End",
     "Diarias premium", "Condo-Hotel Hamam", GOLD),
]
fy = Cm(7.0)
fw = Cm(7.7)
for i, (n, t, p, v, sub, accent) in enumerate(fases30):
    x = Cm(1.5) + i * (fw + Cm(0.2))
    add_rect(s, x, fy, fw, Cm(9.5), fill=NAVY, radius=True)
    add_rect(s, x, fy, fw, Cm(0.5), fill=accent)
    add_text(s, x + Cm(0.5), fy + Cm(0.7), fw - Cm(1), Cm(0.8),
             n, font=FONT_SANS, size=10, color=accent,
             bold=True, letter_spacing=4)
    add_text(s, x + Cm(0.5), fy + Cm(1.6), fw - Cm(1), Cm(1.4),
             t, font=FONT_SERIF, size=18, color=WHITE, bold=True,
             line_spacing=1.15)
    add_line(s, x + Cm(0.5), fy + Cm(3.5), x + fw - Cm(0.5), fy + Cm(3.5),
             color=accent, weight=0.8)
    add_text(s, x + Cm(0.5), fy + Cm(3.7), fw - Cm(1), Cm(0.6),
             "PUBLICO", font=FONT_SANS, size=8, color=accent,
             bold=True, letter_spacing=3)
    add_text(s, x + Cm(0.5), fy + Cm(4.4), fw - Cm(1), Cm(1.5),
             p, font=FONT_SANS, size=11, color=CREAM, line_spacing=1.3)
    add_text(s, x + Cm(0.5), fy + Cm(6.3), fw - Cm(1), Cm(0.6),
             "PRECIFICACAO", font=FONT_SANS, size=8, color=accent,
             bold=True, letter_spacing=3)
    add_text(s, x + Cm(0.5), fy + Cm(7.0), fw - Cm(1), Cm(1.4),
             v, font=FONT_SERIF, size=14, color=WHITE, bold=True,
             line_spacing=1.2)
    add_text(s, x + Cm(0.5), fy + Cm(8.5), fw - Cm(1), Cm(0.8),
             sub, font=FONT_SANS, size=9, color=CREAM, italic=True)
add_footer(s, 30)


# --- S31 CRONOGRAMA EXECUTIVO ---------------------------------------
s = new_slide(VOID)
add_section_header(s, "30  ·  CRONOGRAMA EXECUTIVO",
                   "Marcos de execucao mapeados.")

crono = [
    ("Masterplan KV (em execucao)", "120 dias", TURQ),
    ("Aprovacoes e licencas", "+120 a 240 dias", TURQ),
    ("Lancamento Fase 1", "+6 meses apos aprovacoes", TURQ),
    ("TVO parcial — primeiras entregas", "A partir do lancamento", CREAM),
    ("Execucao completa Fase 1", "3 a 4 anos", CREAM),
    ("Lancamento Fase 2", "2027 - 2028", CREAM),
    ("Expansao continua", "2028+", GREY_50),
]
ty = Cm(6.5)
for i, (m, p, c) in enumerate(crono):
    yy = ty + i * Cm(1.55)
    # marker dot
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                             Cm(2.2), yy + Cm(0.4), Cm(0.6), Cm(0.6))
    dot.fill.solid()
    dot.fill.fore_color.rgb = c
    dot.line.fill.background()
    dot.shadow.inherit = False
    if i < len(crono) - 1:
        add_line(s, Cm(2.5), yy + Cm(1.0), Cm(2.5), yy + Cm(2.05),
                 color=GREY_30, weight=0.75)
    add_rect(s, Cm(3.5), yy + Cm(0.1), Cm(28.8), Cm(1.3), fill=NAVY, radius=True)
    add_text(s, Cm(4.0), yy + Cm(0.4), Cm(20), Cm(0.8),
             m, font=FONT_SERIF, size=14, color=WHITE)
    add_text(s, Cm(24), yy + Cm(0.4), Cm(8), Cm(0.8),
             p, font=FONT_SERIF, size=14, color=c, bold=True,
             align=PP_ALIGN.RIGHT)

add_footer(s, 31)


# --- S32 PREMISSAS FINANCEIRAS --------------------------------------
s = new_slide(VOID)
add_section_header(s, "31  ·  PREMISSAS FINANCEIRAS",
                   "Margem capturada integralmente. Capital 100% equity.")

prem = [
    ("METODOLOGIA", "FipeZAP  ·  ADEMI  ·  ABRAINC  ·  CBIC  ·  comparaveis nacionais"),
    ("MARGEM DO INCORPORADOR",
     "A tese captura integralmente a margem: entrada a preco de custo, "
     "saida a preco de mercado."),
    ("DISTRIBUICAO PROPORCIONAL", "Via cotas da S.A.  ·  registro em cartorio"),
    ("TRANSPARENCIA",
     "Demonstracoes periodicas  ·  auditoria independente"),
    ("AUSENCIA DE ALAVANCAGEM",
     "Capital 100% equity  ·  o tempo joga a favor do ativo  ·  zero "
     "exposicao a juros"),
]
py = Cm(6.5)
for i, (t, d) in enumerate(prem):
    yy = py + i * Cm(2.0)
    add_rect(s, Cm(1.5), yy, Cm(30.8), Cm(1.8), fill=NAVY, radius=True)
    add_rect(s, Cm(1.5), yy, Cm(0.12), Cm(1.8), fill=TURQ)
    add_text(s, Cm(2.2), yy + Cm(0.25), Cm(8), Cm(0.7),
             t, font=FONT_SANS, size=10, color=TURQ, bold=True,
             letter_spacing=3)
    add_text(s, Cm(2.2), yy + Cm(0.95), Cm(29), Cm(0.8),
             d, font=FONT_SERIF, size=12, color=CREAM, line_spacing=1.4)

add_footer(s, 32)


# --- S33 COMPARATIVO INVESTIMENTOS ----------------------------------
s = new_slide(VOID)
add_section_header(s, "32  ·  JUBIABA vs. ALTERNATIVAS",
                   "Real estate + localizacao inexplorada + timing.")

ty = Cm(6.5)
add_rect(s, Cm(1.5), ty, Cm(30.8), Cm(0.85), fill=NAVY_LIGHT)
hdr = [("METRICA", Cm(2)), ("RENDA FIXA", Cm(9)), ("BOLSA", Cm(15.5)),
       ("IMOVEL URBANO", Cm(20.5)), ("JUBIABA F1", Cm(27))]
for h, hx in hdr:
    add_text(s, hx, ty + Cm(0.2), Cm(7), Cm(0.6), h,
             font=FONT_SANS, size=9, color=TURQ, bold=True, letter_spacing=2)

rows = [
    ("Retorno Medio", "10-12% a.a.", "Variavel", "0,5-1% a.m.", "60-120% a.a."),
    ("Risco", "Baixo", "Alto", "Baixo", "Mitigado"),
    ("Liquidez", "Alta", "Alta", "Baixa", "Estruturada"),
    ("Perfil", "Protecao", "Volatil", "Maduro", "Real + Timing"),
]
for i, r in enumerate(rows):
    ry = ty + Cm(0.85) + i * Cm(1.2)
    if i % 2 == 0:
        add_rect(s, Cm(1.5), ry, Cm(30.8), Cm(1.2), fill=NAVY_LIGHT)
    add_text(s, Cm(2), ry + Cm(0.4), Cm(7), Cm(0.7), r[0],
             font=FONT_SERIF, size=12, color=WHITE, bold=True)
    for j, val in enumerate(r[1:]):
        col = TURQ if j == 3 else CREAM
        weight = (j == 3)
        f = FONT_SERIF if j == 3 else FONT_SANS
        add_text(s, hdr[j+1][1], ry + Cm(0.4), Cm(6.5), Cm(0.7), val,
                 font=f, size=11, color=col, bold=weight)

# Formula
add_rect(s, Cm(1.5), Cm(13.5), Cm(30.8), Cm(2.5), fill=NAVY, radius=True)
add_rect(s, Cm(1.5), Cm(13.5), Cm(0.12), Cm(2.5), fill=GOLD)
add_text(s, Cm(2.2), Cm(13.85), Cm(20), Cm(0.7),
         "A FORMULA DO RETORNO",
         font=FONT_SANS, size=10, color=GOLD, bold=True, letter_spacing=4)
add_text(s, Cm(2.2), Cm(14.6), Cm(29), Cm(1.7),
         "Real Estate Tangivel  +  Localizacao Inexplorada  +  Timing (Fase 1)",
         font=FONT_SERIF, size=18, color=WHITE, italic=True, line_spacing=1.3)
add_footer(s, 33)


# --- S34 RETORNO EXPLORER -------------------------------------------
def projection_slide(idx, eyebrow, title, total, cotas, m2, custo,
                     proj, recommend=False, slide_num=0):
    s_ = new_slide(VOID)
    add_section_header(s_, eyebrow, title)
    if recommend:
        # Star badge
        add_rect(s_, Cm(28), Cm(5.5), Cm(4.5), Cm(0.8), fill=GOLD, radius=True)
        add_text(s_, Cm(28), Cm(5.55), Cm(4.5), Cm(0.7),
                 "★ RECOMENDADO", font=FONT_SANS, size=10, color=VOID,
                 bold=True, align=PP_ALIGN.CENTER, letter_spacing=4)

    # Top stats
    stats_p = [
        ("INVESTIMENTO", total),
        ("COTAS", cotas),
        ("AREA EQUIVALENTE", m2),
        ("CUSTO POR M2", custo),
    ]
    for i, (lbl, val) in enumerate(stats_p):
        x = Cm(1.5) + i * (Cm(7.7) + Cm(0.2))
        add_rect(s_, x, Cm(6.5), Cm(7.7), Cm(3.0), fill=NAVY, radius=True)
        add_rect(s_, x, Cm(6.5), Cm(0.12), Cm(3.0), fill=TURQ)
        add_text(s_, x + Cm(0.5), Cm(6.85), Cm(7), Cm(0.6),
                 lbl, font=FONT_SANS, size=9, color=TURQ,
                 bold=True, letter_spacing=3)
        add_text(s_, x + Cm(0.5), Cm(7.55), Cm(7), Cm(2),
                 val, font=FONT_SERIF, size=22, color=WHITE, bold=True,
                 line_spacing=1.1)

    # Cenarios
    cy = Cm(10.3)
    cw = Cm(10.05)
    cenarios_lbl = ["CONSERVADOR", "REALISTA (BASE)", "OTIMISTA"]
    accents_p = [GREY_50, TURQ, GOLD]
    for i, (label, accent) in enumerate(zip(cenarios_lbl, accents_p)):
        x = Cm(1.5) + i * (cw + Cm(0.4))
        add_rect(s_, x, cy, cw, Cm(6.4), fill=NAVY, radius=True)
        add_rect(s_, x, cy, cw, Cm(0.5), fill=accent)
        add_text(s_, x + Cm(0.6), cy + Cm(0.7), cw - Cm(1), Cm(0.7),
                 label, font=FONT_SANS, size=10, color=accent,
                 bold=True, letter_spacing=3)
        add_text(s_, x + Cm(0.6), cy + Cm(1.7), cw - Cm(1), Cm(2.5),
                 proj[i][0], font=FONT_SERIF, size=30, color=WHITE,
                 bold=True, line_spacing=1.05)
        add_text(s_, x + Cm(0.6), cy + Cm(4.4), cw - Cm(1), Cm(1.0),
                 proj[i][1], font=FONT_SERIF, size=22, color=accent,
                 bold=True)
        if len(proj[i]) > 2 and proj[i][2]:
            add_text(s_, x + Cm(0.6), cy + Cm(5.5), cw - Cm(1), Cm(0.7),
                     proj[i][2], font=FONT_SANS, size=10, color=CREAM,
                     italic=True)

    add_disclaimer(s_)
    add_footer(s_, slide_num)


projection_slide(
    34, "33  ·  PROJECAO DE RETORNO",
    "EXPLORER  ·  Diversificacao de portfolio.",
    "R$ 300.000", "1 cota", "54,57 m2", "R$ 5.497 / m2",
    [
        ("R$ 654.840", "2,18x", ""),
        ("R$ 944.143", "3,14x", ""),
        ("R$ 1.098.213", "3,66x", ""),
    ], slide_num=34)


# --- S35 FOUNDER ----------------------------------------------------
projection_slide(
    35, "34  ·  PROJECAO DE RETORNO",
    "FOUNDER  ·  O ponto de equilibrio perfeito.",
    "R$ 1.200.000", "4 cotas", "266,80 m2", "R$ 4.497 / m2",
    [
        ("R$ 3.201.600", "2,67x", "R$ 12.000 / m2"),
        ("R$ 4.613.440", "3,84x", "R$ 17.291 / m2"),
        ("R$ 5.360.680", "4,47x", "R$ 20.092 / m2"),
    ], recommend=True, slide_num=35)


# --- S36 LEGACY -----------------------------------------------------
projection_slide(
    36, "35  ·  PROJECAO DE RETORNO",
    "LEGACY  ·  Family Office  ·  Institucional.",
    "R$ 2.100.000+", "7+ cotas", "525 m2+", "R$ 4.000 / m2",
    [
        ("R$ 6.300.000", "3,00x", "R$ 12.000 / m2"),
        ("R$ 9.082.500", "4,33x", "R$ 17.291 / m2"),
        ("R$ 10.552.500", "5,03x", "R$ 20.092 / m2"),
    ], slide_num=36)


# --- S37 ANALISE DE SENSIBILIDADE -----------------------------------
s = new_slide(VOID)
add_section_header(s, "36  ·  ANALISE DE SENSIBILIDADE",
                   "Tres cenarios. Tres testes de estresse.")

# Cenarios row
ce = [
    ("CONSERVADOR", "R$ 12.000 / m2",
     "Valorizacao estagnada — sem prêmio de infraestrutura ou marca.\n"
     "Mesmo aqui: 2,18x minimo.", GREY_50),
    ("REALISTA (BASE)", "R$ 17.291 / m2",
     "Correcao inflacionaria + premio de desenvolvimento\n"
     "padrao resort de luxo.", TURQ),
    ("OTIMISTA", "R$ 20.092 / m2",
     "Alta demanda consolidando Jubiaba como\n"
     "novo polo nacional — Efeito Jurere/Trancoso.", GOLD),
]
for i, (lbl, val, desc, c) in enumerate(ce):
    x = Cm(1.5) + i * Cm(10.5)
    add_rect(s, x, Cm(6.5), Cm(10.2), Cm(5.5), fill=NAVY, radius=True)
    add_rect(s, x, Cm(6.5), Cm(10.2), Cm(0.4), fill=c)
    add_text(s, x + Cm(0.5), Cm(7.0), Cm(9), Cm(0.7),
             lbl, font=FONT_SANS, size=10, color=c,
             bold=True, letter_spacing=3)
    add_text(s, x + Cm(0.5), Cm(7.8), Cm(9), Cm(1.4),
             val, font=FONT_SERIF, size=22, color=WHITE, bold=True)
    add_text(s, x + Cm(0.5), Cm(9.6), Cm(9), Cm(2.4),
             desc, font=FONT_SANS, size=10, color=CREAM, line_spacing=1.4)

# Stress test
add_text(s, Cm(1.5), Cm(12.5), Cm(20), Cm(0.7),
         "TESTE DE ESTRESSE  ·  3 VARIAVEIS",
         font=FONT_SANS, size=10, color=TURQ, bold=True, letter_spacing=4)
stress = [
    ("PRECO DE VENDA", "+/- 10%", "ALTO", "principal motor da margem"),
    ("CRONOGRAMA", "+/- 6 meses", "MEDIO", "sem divida bancaria, nao corroi capital"),
    ("TAXA DE OCUPACAO", "+/- 20%", "MEDIO", "afeta receita hoteleira, nao o VGV"),
]
for i, (var, delta, imp, obs) in enumerate(stress):
    yy = Cm(13.5) + i * Cm(1.2)
    add_rect(s, Cm(1.5), yy, Cm(30.8), Cm(1.1), fill=NAVY, radius=True)
    add_text(s, Cm(2), yy + Cm(0.3), Cm(8), Cm(0.7),
             var, font=FONT_SERIF, size=12, color=WHITE, bold=True)
    add_text(s, Cm(10), yy + Cm(0.32), Cm(5), Cm(0.7),
             delta, font=FONT_SANS, size=11, color=CREAM)
    color_imp = TURQ if imp == "ALTO" else GOLD
    add_text(s, Cm(15), yy + Cm(0.3), Cm(5), Cm(0.7),
             imp, font=FONT_SANS, size=11, color=color_imp,
             bold=True, letter_spacing=3)
    add_text(s, Cm(20), yy + Cm(0.32), Cm(12), Cm(0.7),
             obs, font=FONT_SANS, size=10, color=CREAM, italic=True)

add_disclaimer(s)
add_footer(s, 37)


# --- S38 PAYBACK ----------------------------------------------------
s = new_slide(VOID)
add_section_header(s, "37  ·  PAYBACK DO INVESTIDOR",
                   "Da assinatura ao lucro puro.")

milestones = [
    ("M0", "APORTE", "Capital aportado\nGarantia Real imediata", TURQ),
    ("M36-42", "PRIMEIRA RECEITA",
     "Inicio das vendas Etapa 01\n(Seaside)\nPrimeiras distribuicoes", CREAM),
    ("M48-56", "PAYBACK",
     "Retorno integral do capital\nRisco financeiro zerado", GOLD),
    ("M56+", "LUCRO PURO",
     "Distribuicao dos lucros\n+\nReceita hoteleira recorrente Hamam", TURQ),
]
my = Cm(6.5)
mw = Cm(7.7)
for i, (m, t, d, c) in enumerate(milestones):
    x = Cm(1.5) + i * (mw + Cm(0.2))
    add_rect(s, x, my, mw, Cm(8), fill=NAVY, radius=True)
    add_rect(s, x, my, mw, Cm(0.4), fill=c)
    add_text(s, x + Cm(0.5), my + Cm(0.6), mw - Cm(1), Cm(1.5),
             m, font=FONT_SERIF, size=24, color=c, bold=True)
    add_text(s, x + Cm(0.5), my + Cm(2.4), mw - Cm(1), Cm(1),
             t, font=FONT_SANS, size=10, color=WHITE,
             bold=True, letter_spacing=3)
    add_line(s, x + Cm(0.5), my + Cm(3.5), x + Cm(2.5), my + Cm(3.5),
             color=c, weight=1.0)
    add_text(s, x + Cm(0.5), my + Cm(3.9), mw - Cm(1), Cm(4),
             d, font=FONT_SANS, size=10, color=CREAM, line_spacing=1.5)

# Exemplo box
add_rect(s, Cm(1.5), Cm(15.0), Cm(30.8), Cm(2.5), fill=NAVY_LIGHT, radius=True)
add_rect(s, Cm(1.5), Cm(15.0), Cm(0.12), Cm(2.5), fill=GOLD)
add_text(s, Cm(2.2), Cm(15.3), Cm(15), Cm(0.6),
         "EXEMPLO PRATICO  ·  COTA EXPLORER",
         font=FONT_SANS, size=10, color=GOLD, bold=True, letter_spacing=4)
add_text(s, Cm(2.2), Cm(16.0), Cm(29), Cm(1.5),
         "R$ 300.000  →  Lucro de R$ 354k a R$ 798k  ·  Multiplo 2,18x a 3,66x",
         font=FONT_SERIF, size=16, color=WHITE, italic=True)
add_disclaimer(s)
add_footer(s, 38)


# --- S39 ESTRUTURA DE SAIDA: 3 OPCOES -------------------------------
s = new_slide(VOID)
add_section_header(s, "38  ·  ESTRUTURA DE SAIDA",
                   "Tres opcoes. Maxima flexibilidade.")

opcoes = [
    ("01", "VENDA PARA LUCRO",
     "Revenda a mercado.\nRealizacao do multiplicador 3x - 6x.\n\n"
     "Foco em liquidez maxima.", TURQ),
    ("02", "PERMANENCIA",
     "Uso proprio (segunda residencia)\nou locacao tradicional.\n\n"
     "Preservacao de patrimonio em ativo real.", CREAM),
    ("03", "CONDO-HOTEL",
     "Pool hoteleiro gerido pela Hamam Global.\nReceita passiva recorrente.\n\n"
     "Gestao profissional end-to-end.", GOLD),
]
oy = Cm(6.5)
ow = Cm(10.2)
for i, (n, t, d, c) in enumerate(opcoes):
    x = Cm(1.5) + i * (ow + Cm(0.4))
    add_rect(s, x, oy, ow, Cm(11.0), fill=NAVY, radius=True)
    add_text(s, x + Cm(0.7), oy + Cm(0.7), Cm(3), Cm(2.5),
             n, font=FONT_SERIF, size=64, color=c, line_spacing=1.0)
    add_line(s, x + Cm(0.7), oy + Cm(3.7), x + ow - Cm(0.7), oy + Cm(3.7),
             color=c, weight=1.0)
    add_text(s, x + Cm(0.7), oy + Cm(4.0), ow - Cm(1.4), Cm(1.5),
             t, font=FONT_SERIF, size=20, color=WHITE, bold=True,
             line_spacing=1.15)
    add_text(s, x + Cm(0.7), oy + Cm(6.5), ow - Cm(1.4), Cm(4),
             d, font=FONT_SANS, size=12, color=CREAM, line_spacing=1.5)

add_footer(s, 39)


# --- S40 USO DOS RECURSOS ROUND 02 ----------------------------------
s = new_slide(VOID)
add_section_header(s, "39  ·  USO DOS RECURSOS  ·  ROUND 02",
                   "R$ 10 milhoes  ·  destinacao carimbada.")

usos = [
    ("ESTRUTURACAO & APROVACOES",
     "Projetos executivos KV  ·  licencas  ·  taxas  ·  SPE", "30%"),
    ("INFRA & ACESSO",
     "Terraplenagem  ·  pavimentacao  ·  redes primarias Etapa 01", "35%"),
    ("CONSTRUCAO ETAPA 01",
     "Fundacoes  ·  primeiras unidades Seaside  ·  areas comuns", "25%"),
    ("FUNDO DE CONTINGENCIA",
     "Reserva para variacoes de custo e imprevistos operacionais", "10%"),
]
uy = Cm(6.8)
uw = Cm(15.4)
for i, (t, d, p) in enumerate(usos):
    col = i % 2
    row = i // 2
    x = Cm(1.5) + col * (uw + Cm(0.3))
    y = uy + row * Cm(5.4)
    add_rect(s, x, y, uw, Cm(5.1), fill=NAVY, radius=True)
    add_text(s, x + Cm(0.6), y + Cm(0.4), Cm(8), Cm(0.7),
             t, font=FONT_SANS, size=10, color=TURQ,
             bold=True, letter_spacing=3)
    add_text(s, x + uw - Cm(4.5), y + Cm(0.2), Cm(4), Cm(2.2),
             p, font=FONT_SERIF, size=44, color=TURQ, bold=True,
             align=PP_ALIGN.RIGHT, line_spacing=1.0)
    add_text(s, x + Cm(0.6), y + Cm(2.5), uw - Cm(1), Cm(2.5),
             d, font=FONT_SERIF, size=13, color=CREAM,
             italic=True, line_spacing=1.4)

add_footer(s, 40)


# ======================================================================
# BLOCO 05 - GOVERNANCA, SEGURANCA & EQUIPE
# ======================================================================

# --- S41 VANTAGENS COMPETITIVAS / MOAT ------------------------------
s = new_slide(VOID)
add_section_header(s, "40  ·  O MOAT DO PROJETO",
                   "Seis diferenciais estruturais  ·  barreira intransponivel.")

moat = [
    ("LITORAL VIRGEM",
     "Praias intocadas  ·  privacidade  ·  contraste com destinos saturados"),
    ("SMART CITY PLANTA NOVA",
     "Infraestrutura tecnologica subterranea desde o zero"),
    ("HOSPITALIDADE INTEGRADA",
     "Hamam Global  ·  padrao internacional desde o primeiro dia"),
    ("ARQUITETURA DE AUTOR",
     "KV  ·  valor estetico inquestionavel  ·  chancela premium"),
    ("SUSTENTABILIDADE REAL",
     "ZEEC aprovado  ·  dark-sky  ·  APP como parque natural"),
    ("ACESSIBILIDADE ESTRATEGICA",
     "50 min do aeroporto  ·  Linha Verde  ·  polo industrial"),
]
gx = Cm(1.5)
gy = Cm(6.5)
gw = Cm(10.05)
gh = Cm(5.4)
for i, (t, d) in enumerate(moat):
    col = i % 3
    row = i // 3
    x = gx + col * (gw + Cm(0.4))
    y = gy + row * (gh + Cm(0.4))
    add_rect(s, x, y, gw, gh, fill=NAVY, radius=True)
    add_text(s, x + Cm(0.7), y + Cm(0.5), Cm(2), Cm(2),
             f"0{i+1}", font=FONT_SERIF, size=42, color=TURQ, line_spacing=1.0)
    add_line(s, x + Cm(0.7), y + Cm(3.0), x + Cm(2.8), y + Cm(3.0),
             color=TURQ, weight=1.0)
    add_text(s, x + Cm(0.7), y + Cm(3.2), gw - Cm(1), Cm(0.8),
             t, font=FONT_SANS, size=10, color=WHITE,
             bold=True, letter_spacing=3)
    add_text(s, x + Cm(0.7), y + Cm(4.1), gw - Cm(1), Cm(1.4),
             d, font=FONT_SANS, size=10, color=CREAM, line_spacing=1.4)
add_footer(s, 41)


# --- S42 MATRIZ DE RISCOS -------------------------------------------
s = new_slide(VOID)
add_section_header(s, "41  ·  MATRIZ DE RISCOS",
                   "Cinco riscos. Todos mapeados antes da captacao.")

riscos = [
    ("MERCADO", "MITIGADO",
     "Demanda comprovada (9x)  ·  escassez de produtos premium", TURQ),
    ("EXECUCAO", "MITIGADO",
     "Sistema GSWP  ·  parcerias KV e Hamam  ·  cronograma estruturado", TURQ),
    ("AMBIENTAL", "MITIGADO",
     "ZEEC aprovado  ·  compliance-first desde a concepcao", TURQ),
    ("REGULATORIO", "MITIGADO",
     "Aprovacoes municipais regularizadas  ·  sem passivos ocultos", TURQ),
    ("FINANCEIRO", "BLINDADO",
     "Aporte 100% lastreado em garantia real  ·  registro em cartorio", GOLD),
]
ry = Cm(6.5)
for i, (t, status, d, c) in enumerate(riscos):
    yy = ry + i * Cm(2.1)
    add_rect(s, Cm(1.5), yy, Cm(30.8), Cm(1.9), fill=NAVY, radius=True)
    add_rect(s, Cm(1.5), yy, Cm(0.12), Cm(1.9), fill=c)
    add_text(s, Cm(2.2), yy + Cm(0.3), Cm(7), Cm(0.7),
             t, font=FONT_SERIF, size=18, color=WHITE, bold=True)
    # Status badge
    add_rect(s, Cm(2.2), yy + Cm(1.05), Cm(3), Cm(0.65), fill=c)
    add_text(s, Cm(2.2), yy + Cm(1.1), Cm(3), Cm(0.6),
             status, font=FONT_SANS, size=9, color=VOID,
             bold=True, align=PP_ALIGN.CENTER, letter_spacing=3)
    add_text(s, Cm(10), yy + Cm(0.55), Cm(22), Cm(1.2),
             d, font=FONT_SERIF, size=13, color=CREAM, italic=True,
             line_spacing=1.4)
add_footer(s, 42)


# --- S43 ESG & IMPACTO ----------------------------------------------
s = new_slide(VOID)
add_section_header(s, "42  ·  ESG & IMPACTO REGIONAL",
                   "Sustentabilidade nao como narrativa  ·  como concepcao.")

esg = [
    ("28-35%", "APP como\nParque Natural integrado",
     "Cordao dunar preservado em sua integridade."),
    ("Dark-Sky", "Protecao a\ntartarugas marinhas",
     "Diferencial de consciencia ambiental e marketing."),
    ("Solar", "Energia limpa\n+ SUDS + ETE",
     "Sem emissao na orla. Matriz totalmente renovavel."),
    ("Local", "Geracao de\nemprego e renda",
     "Cadeia produtiva regional fortalecida."),
    ("Ciclovia", "Mobilidade\nlimpa integrada",
     "Vias com ciclofaixas  ·  shuttle eletrico interno."),
]
gx = Cm(1.5)
gy = Cm(6.5)
gw = Cm(6.05)
gh = Cm(11)
for i, (val, t, d) in enumerate(esg):
    x = gx + i * (gw + Cm(0.2))
    add_rect(s, x, gy, gw, gh, fill=NAVY, radius=True)
    add_rect(s, x, gy, gw, Cm(0.5), fill=TURQ)
    add_text(s, x + Cm(0.4), gy + Cm(0.9), gw - Cm(0.8), Cm(2),
             val, font=FONT_SERIF, size=22, color=TURQ, bold=True,
             align=PP_ALIGN.CENTER, line_spacing=1.0)
    add_line(s, x + Cm(1.5), gy + Cm(3.4), x + gw - Cm(1.5), gy + Cm(3.4),
             color=TURQ, weight=0.8)
    add_text(s, x + Cm(0.4), gy + Cm(3.8), gw - Cm(0.8), Cm(2.5),
             t, font=FONT_SERIF, size=14, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, line_spacing=1.2)
    add_text(s, x + Cm(0.4), gy + Cm(7.0), gw - Cm(0.8), Cm(3.5),
             d, font=FONT_SANS, size=10, color=CREAM,
             align=PP_ALIGN.CENTER, line_spacing=1.45)
add_footer(s, 43)


# --- S44 EQUIPE -----------------------------------------------------
s = new_slide(VOID)
add_section_header(s, "43  ·  A EQUIPE",
                   "Tres lideres. Tres operacoes de classe mundial.")

equipe = [
    ("BRUNO ZANUTO",
     "CEO & Fundador  ·  GSWP",
     "Desenvolvimento imobiliario e governanca.\n"
     "Operacao validada por fundos institucionais.\n"
     "Sistema construtivo proprietario."),
    ("KONIGSBERGER VANNUCCHI",
     "Masterplan & Arquitetura",
     "53 anos de operacao.\nReferencia absoluta do Brasil em masterplans "
     "+ residencial premium.\n893 projetos  ·  R$ 18Bi de VGV."),
    ("HAMAM GLOBAL",
     "Hospitalidade integrada",
     "26 anos de operacao.\nLider da America Latina em FF&E/OS&E.\n"
     "3 hubs globais  ·  +250 hoteis entregues."),
]
ey = Cm(6.5)
ew = Cm(10.2)
for i, (nome, papel, desc) in enumerate(equipe):
    x = Cm(1.5) + i * (ew + Cm(0.4))
    add_rect(s, x, ey, ew, Cm(11.0), fill=NAVY, radius=True)
    # placeholder portrait
    add_rect(s, x + Cm(0.5), ey + Cm(0.5), ew - Cm(1), Cm(4.5),
             fill=NAVY_LIGHT, radius=True)
    # initial badge
    initials = "BZ" if i == 0 else ("KV" if i == 1 else "HG")
    add_text(s, x + Cm(0.5), ey + Cm(1.4), ew - Cm(1), Cm(3),
             initials, font=FONT_SERIF, size=64, color=TURQ, bold=True,
             align=PP_ALIGN.CENTER, line_spacing=1.0)
    add_line(s, x + Cm(0.7), ey + Cm(5.5), x + ew - Cm(0.7), ey + Cm(5.5),
             color=TURQ, weight=1.0)
    add_text(s, x + Cm(0.7), ey + Cm(5.7), ew - Cm(1.4), Cm(1.2),
             nome, font=FONT_SERIF, size=18, color=WHITE, bold=True)
    add_text(s, x + Cm(0.7), ey + Cm(6.9), ew - Cm(1.4), Cm(0.8),
             papel.upper(), font=FONT_SANS, size=9, color=TURQ,
             bold=True, letter_spacing=3)
    add_text(s, x + Cm(0.7), ey + Cm(7.9), ew - Cm(1.4), Cm(3),
             desc, font=FONT_SANS, size=10, color=CREAM, line_spacing=1.45)
add_footer(s, 44)


# ======================================================================
# BLOCO 06 - CALL TO ACTION
# ======================================================================

# --- S45 ESCOLHA SUA PARTICIPACAO -----------------------------------
s = new_slide(VOID)
add_section_header(s, "44  ·  ESCOLHA SUA PARTICIPACAO",
                   "Tres tickets. Maximizando poder de compra.")

tickets = [
    ("EXPLORER", "1 cota", "R$ 300.000", "54,57 m2", "R$ 5.497 / m2",
     "Diversificacao\nde portfolio", GREY_50, False),
    ("FOUNDER", "4 cotas", "R$ 1.200.000", "266,80 m2", "R$ 4.497 / m2",
     "Ponto de equilibrio\nperfeito", TURQ, True),
    ("LEGACY", "7+ cotas", "R$ 2.100.000+", "525+ m2", "R$ 4.000 / m2",
     "Family Office\nInstitucional", GOLD, False),
]
tx = Cm(1.5)
ty = Cm(6.5)
tw = Cm(10.3)
for i, (nome, cotas, inv, area, custo, perfil, c, rec) in enumerate(tickets):
    x = tx + i * (tw + Cm(0.35))
    fill = TURQ if rec else NAVY
    add_rect(s, x, ty, tw, Cm(11.0), fill=fill, radius=True)
    txt_main = VOID if rec else WHITE
    txt_sec = VOID if rec else c
    txt_body = VOID if rec else CREAM

    if rec:
        add_text(s, x + Cm(0.6), ty + Cm(0.5), tw - Cm(1), Cm(0.7),
                 "★ RECOMENDADO", font=FONT_SANS, size=10, color=VOID,
                 bold=True, letter_spacing=4)

    yhead = ty + (Cm(1.4) if rec else Cm(0.7))
    add_text(s, x + Cm(0.6), yhead, tw - Cm(1), Cm(1.5),
             nome, font=FONT_SERIF, size=28, color=txt_main, bold=True)
    add_text(s, x + Cm(0.6), yhead + Cm(1.6), tw - Cm(1), Cm(0.7),
             cotas.upper(), font=FONT_SANS, size=10, color=txt_sec,
             bold=True, letter_spacing=3)
    add_line(s, x + Cm(0.6), yhead + Cm(2.6), x + tw - Cm(0.6), yhead + Cm(2.6),
             color=txt_sec, weight=1.0)

    fields = [
        ("INVESTIMENTO", inv),
        ("AREA EQUIVALENTE", area),
        ("CUSTO POR M2", custo),
    ]
    for j, (lbl, val) in enumerate(fields):
        yyy = yhead + Cm(2.9) + j * Cm(1.55)
        add_text(s, x + Cm(0.6), yyy, tw - Cm(1), Cm(0.5),
                 lbl, font=FONT_SANS, size=8, color=txt_sec,
                 bold=True, letter_spacing=3)
        add_text(s, x + Cm(0.6), yyy + Cm(0.45), tw - Cm(1), Cm(1),
                 val, font=FONT_SERIF, size=15, color=txt_main, bold=True)

    add_text(s, x + Cm(0.6), ty + Cm(9.1), tw - Cm(1), Cm(1.5),
             perfil, font=FONT_SERIF, size=12, color=txt_main, italic=True,
             line_spacing=1.3)

add_text(s, Cm(1.5), Cm(17.7), Cm(30.8), Cm(0.5),
         "Os valores de m2 sao progressivamente menores para maior volume "
         "— maximizando poder de compra.",
         font=FONT_SERIF, size=10, color=GOLD, italic=True,
         align=PP_ALIGN.CENTER)
add_disclaimer(s)
add_footer(s, 45)


# --- S46 PROXIMOS PASSOS --------------------------------------------
s = new_slide(VOID)
add_section_header(s, "45  ·  PROXIMOS PASSOS",
                   "Processo de adesao objetivo.")

passos = [
    ("01", "INTERESSE", "Contato com GSWP"),
    ("02", "NDA + DATA ROOM", "Material completo confidencial"),
    ("03", "DUE DILIGENCE",
     "Acesso ao projeto executivo e juridico"),
    ("04", "CONTRATO", "Contrato de Investimento estruturado"),
    ("05", "APORTE", "Capital com Garantia Real ativa"),
]
py = Cm(6.5)
pw = Cm(6.05)
for i, (n, t, d) in enumerate(passos):
    x = Cm(1.5) + i * (pw + Cm(0.2))
    add_rect(s, x, py, pw, Cm(7.5), fill=NAVY, radius=True)
    # Number circle
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Cm(0.5), py + Cm(0.5),
                              Cm(2.0), Cm(2.0))
    circ.fill.solid()
    circ.fill.fore_color.rgb = TURQ
    circ.line.fill.background()
    circ.shadow.inherit = False
    add_text(s, x + Cm(0.5), py + Cm(0.65), Cm(2), Cm(1.7),
             n, font=FONT_SERIF, size=22, color=VOID, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, x + Cm(0.4), py + Cm(2.9), pw - Cm(0.8), Cm(1.4),
             t, font=FONT_SANS, size=10, color=WHITE,
             bold=True, letter_spacing=2)
    add_line(s, x + Cm(0.5), py + Cm(4.3), x + pw - Cm(0.5), py + Cm(4.3),
             color=TURQ, weight=0.8)
    add_text(s, x + Cm(0.4), py + Cm(4.6), pw - Cm(0.8), Cm(2.5),
             d, font=FONT_SERIF, size=11, color=CREAM,
             italic=True, line_spacing=1.45)

# Round 02 highlight
add_rect(s, Cm(1.5), Cm(15.0), Cm(30.8), Cm(2.5), fill=TURQ, radius=True)
add_text(s, Cm(2.2), Cm(15.3), Cm(15), Cm(0.6),
         "ROUND 02", font=FONT_SANS, size=10, color=VOID,
         bold=True, letter_spacing=4)
add_text(s, Cm(2.2), Cm(16.0), Cm(29), Cm(1.4),
         "R$ 10 milhoes  ·  janela exclusiva como socio fundador da Fase 1",
         font=FONT_SERIF, size=18, color=VOID, bold=True, italic=True)
add_footer(s, 46)


# --- S47 ENCERRAMENTO -----------------------------------------------
s = new_slide(VOID)
if img_aerea2 or img_praia_saco:
    add_image_full(s, img_aerea2 or img_praia_saco)
add_overlay(s, opacity=0.65, color=VOID)

# Eyebrow
add_text(s, Cm(1.5), Cm(2.0), Cm(20), Cm(0.6),
         "JUBIABA  ·  2025 / 2026",
         font=FONT_SANS, size=9, color=TURQ, bold=True, letter_spacing=6)
add_line(s, Cm(1.5), Cm(2.7), Cm(7), Cm(2.7), color=TURQ, weight=1.2)

# Big closing
add_text(s, Cm(1.5), Cm(6.0), Cm(31), Cm(3),
         "OBRIGADO.",
         font=FONT_SERIF, size=88, color=WHITE, bold=False,
         letter_spacing=15, line_spacing=1.0)
add_text(s, Cm(1.5), Cm(9.5), Cm(31), Cm(2),
         "A proxima fronteira do luxo no Nordeste\ncomeca agora.",
         font=FONT_SERIF, size=24, color=CREAM, italic=True, line_spacing=1.4)

# Partners
add_line(s, Cm(1.5), Cm(13.3), Cm(32.3), Cm(13.3), color=TURQ, weight=0.75)
add_text(s, Cm(1.5), Cm(13.5), Cm(31), Cm(0.8),
         "GSWP    ×    KONIGSBERGER VANNUCCHI    ×    HAMAM GLOBAL",
         font=FONT_SANS, size=10, color=WHITE, bold=True,
         letter_spacing=8, align=PP_ALIGN.CENTER)

# Contact
add_rect(s, Cm(1.5), Cm(15.2), Cm(30.8), Cm(2.5), fill=NAVY, radius=True)
add_rect(s, Cm(1.5), Cm(15.2), Cm(0.12), Cm(2.5), fill=TURQ)
add_text(s, Cm(2.2), Cm(15.5), Cm(10), Cm(0.6),
         "CONTATO  ·  ROUND 02",
         font=FONT_SANS, size=10, color=TURQ, bold=True, letter_spacing=4)
add_text(s, Cm(2.2), Cm(16.2), Cm(29), Cm(1.4),
         "investimento@gswp.com.br   ·   www.gswp.com.br   ·   "
         "Bruno Zanuto  ·  CEO",
         font=FONT_SERIF, size=14, color=CREAM, italic=True)

add_text(s, Cm(1.5), SLIDE_H - Cm(0.55), SLIDE_W - Cm(3), Cm(0.4),
         "Material informativo confidencial. Projecoes baseadas em analise "
         "de mercado e historico de valorizacao. Retornos passados nao "
         "garantem resultados futuros. Consulte o prospecto completo.",
         font=FONT_SANS, size=6, color=GREY_30, italic=True)


# ----------------------------------------------------------------------
# Salvar
# ----------------------------------------------------------------------
prs.save(OUT)
print(f"OK -> {OUT}")
print(f"Total slides: {len(prs.slides)}")
