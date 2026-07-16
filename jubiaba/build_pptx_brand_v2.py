# -*- coding: utf-8 -*-
"""
Gerador PPTX v2 — Jubiabá com Brand System GSWP + FOTOS.
Lê _slides_parsed.json (89 slides) e gera jubiaba-gswp-brand-v2.pptx.
Baseado em build_pptx_brand.py, adicionando fotos (background / half-slide)
em slides selecionados, com overlay escuro p/ legibilidade.
"""
import re, sys, json, os
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

try:
    from PIL import Image
    HAS_PIL = True
except Exception:
    HAS_PIL = False

BASE = 'D:/Copy EI/jubiaba'
AREA = os.path.join(BASE, 'FOTOS', 'Fotos da Área')
REGIAO = os.path.join(BASE, 'FOTOS', 'Fotos da Região_')

# ---------------- BRAND SYSTEM GSWP ----------------
VOID     = RGBColor(0x09,0x09,0x07)
CHARCOAL = RGBColor(0x1A,0x19,0x17)
STONE    = RGBColor(0x2C,0x2B,0x28)
GOLD     = RGBColor(0xB8,0x94,0x2E)
GOLD_L   = RGBColor(0xC4,0xA8,0x4A)
CREAM    = RGBColor(0xF0,0xED,0xE6)
OFFWHITE = RGBColor(0xFA,0xF8,0xF4)
SAGE     = RGBColor(0x2D,0x4A,0x3E)
MUTED    = RGBColor(0x9A,0x94,0x90)

# ---------------- FONTES ----------------
SERIF = "Cormorant Garamond"
SANS  = "DM Sans"

# ---------------- SLIDE GEOMETRY (16:9, 1280x720 px -> EMU) ----------------
PX = 9525
SLIDE_W = 1280 * PX
SLIDE_H = 720 * PX

FOOTER = "GSWP × KV × HAMAM GLOBAL · Material Confidencial"

# ---------------- FOTOS por slide ----------------
# tipo: 'bg' (background full) ou 'side' (half-slide direito)
SLIDE_PHOTOS = {
    '01': ('bg',   os.path.join(AREA,   'PHOTO-2023-12-04-09-13-53.jpg'),                         0.62),
    '13': ('bg',   os.path.join(REGIAO, 'praia-do-saco-um-paraiso-proximo-de-aracaju.jpg'),       0.55),
    '17': ('bg',   os.path.join(REGIAO, 'por-do-sol-na-praia-do.jpg'),                            0.58),
    '21': ('side', os.path.join(AREA,   'AEROPORTO ARACAJU.jpeg'),                                 0.30),
    '22': ('bg',   os.path.join(AREA,   'WhatsApp Image 2025-11-14 at 13.14.07.jpeg'),            0.60),
    '23': ('bg',   os.path.join(AREA,   'WhatsApp Image 2025-11-14 at 13.14.05.jpeg'),            0.60),
    '24': ('side', os.path.join(REGIAO, 'Mangue-Seco-e-Praia-do-Saco.jpg'),                       0.30),
    '89': ('bg',   os.path.join(REGIAO, 'praia-do-saco-um-paraiso-proximo-de-aracaju-2.jpg'),     0.55),
}

# ---------------- REFERÊNCIAS por slide (mapa) ----------------
REFS = {
    '02': "FINANCIAL TIMES — Londres perde população",
    '03': "BofA / THE RIO TIMES — Brazil: The New Gold",
    '04': "BofA / THE RIO TIMES — Brazil: The New Gold",
    '05': "CBIC — Luxo no NE +85%",
    '06': "CBIC — Luxo no NE +85%",
    '07': "DESENVOLVE-SE — Turismo como nova fronteira",
    '09': "GOVERNO SE — Aeroporto em volume recorde",
    '11': "PREFEITURA DE ESTÂNCIA — 45% das exportações de SE",
    '12': "REUTERS / PETROBRAS — Sergipe Águas Profundas",
    '21': "GOVERNO SE — Aeroporto em volume recorde",
    '67': "FINANCIAL TIMES — Climate & lifestyle migration",
    '68': "BofA / THE RIO TIMES — Brazil: The New Gold",
    '69': "PODER360 — Sergipe +315% em gás natural",
}

STATS_CAPA = [("250 ha","ÁREA TOTAL"),("1.070 m","FRENTE-MAR"),("R$ 2 bi","VGV")]

TICKETS = [
    ("EXPLORER","R$ 300 mil","54,57 m²"),
    ("FOUNDER ★","R$ 1,2 mi","266,80 m²"),
    ("LEGACY","R$ 2,1 mi +","525 m² +"),
]

# ---------------- HELPERS ----------------
def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if not line:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def set_alpha(shape, pct):
    """Aplica transparência (pct = opacidade 0-100) ao preenchimento sólido."""
    sp = shape.fill._xPr.find(qn('a:solidFill'))
    if sp is None: return
    srgb = sp.find(qn('a:srgbClr'))
    if srgb is None: return
    alpha = srgb.makeelement(qn('a:alpha'), {'val': str(int(pct*1000))})
    srgb.append(alpha)

def _send_to_back(slide, shape, index=2):
    """Move shape para o início do spTree (atrás dos demais)."""
    spTree = slide.shapes._spTree
    spTree.remove(shape._element)
    spTree.insert(index, shape._element)

def add_bg_image(slide, img_path, alpha_overlay=0.55):
    """Adiciona imagem cobrindo o slide inteiro, com overlay escuro semitransparente."""
    if not os.path.isfile(img_path):
        print(f"  [skip bg] arquivo não encontrado: {img_path}")
        return False
    try:
        pic = slide.shapes.add_picture(img_path, Emu(0), Emu(0), Emu(SLIDE_W), Emu(SLIDE_H))
    except Exception as e:
        print(f"  [skip bg] erro ao inserir {os.path.basename(img_path)}: {e}")
        return False
    _send_to_back(slide, pic, 2)
    # overlay escuro sobre a imagem (opacidade da camada escura = alpha_overlay)
    overlay = add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, VOID)
    set_alpha(overlay, int(alpha_overlay * 100))
    _send_to_back(slide, overlay, 3)
    return True

def add_side_image(slide, img_path, alpha_edge=0.30):
    """Adiciona imagem na metade direita (45%) com gradiente/borda escura à esquerda."""
    if not os.path.isfile(img_path):
        print(f"  [skip side] arquivo não encontrado: {img_path}")
        return False
    half_w = int(SLIDE_W * 0.45)
    x0 = SLIDE_W - half_w
    try:
        pic = slide.shapes.add_picture(img_path, Emu(x0), Emu(0), Emu(half_w), Emu(SLIDE_H))
    except Exception as e:
        print(f"  [skip side] erro ao inserir {os.path.basename(img_path)}: {e}")
        return False
    _send_to_back(slide, pic, 2)
    # leve escurecimento geral da imagem p/ integrar à paleta
    tint = add_rect(slide, x0, 0, half_w, SLIDE_H, VOID)
    set_alpha(tint, int(alpha_edge * 100))
    _send_to_back(slide, tint, 3)
    # faixa de transição escura na borda esquerda da foto
    edge = add_rect(slide, x0, 0, int(half_w * 0.32), SLIDE_H, VOID)
    set_alpha(edge, 55)
    _send_to_back(slide, edge, 4)
    return True

def apply_photo(slide, num):
    """Aplica foto configurada para o slide (se houver). Retorna 'bg'|'side'|None."""
    cfg = SLIDE_PHOTOS.get(num)
    if not cfg:
        return None
    kind, path, alpha = cfg
    if kind == 'bg':
        ok = add_bg_image(slide, path, alpha)
        return 'bg' if ok else None
    elif kind == 'side':
        ok = add_side_image(slide, path, alpha)
        return 'side' if ok else None
    return None

def txt(slide, x, y, w, h, text, *, font=SANS, size=14, color=CREAM,
        bold=False, italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        spacing=None, line_spacing=1.1, light=False):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font; f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color
    if spacing is not None:
        rPr = r._r.get_or_add_rPr()
        rPr.set('spc', str(int(spacing*100)))
    return tb

def add_bullets(slide, x, y, w, h, items, *, size=14, color=CREAM, gold_bullet=True):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = 1.15
        p.space_after = Pt(6)
        rb = p.add_run(); rb.text = "◆  "
        rb.font.name = SANS; rb.font.size = Pt(size-3); rb.font.color.rgb = GOLD
        r = p.add_run(); r.text = it
        r.font.name = SANS; r.font.size = Pt(size); r.font.color.rgb = color
    return tb

def diamond(slide, cx, y, color=GOLD):
    size = 10*PX
    d = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Emu(int(cx-size/2)), Emu(y), Emu(size), Emu(size))
    d.fill.solid(); d.fill.fore_color.rgb = color
    d.line.fill.background(); d.shadow.inherit = False
    set_alpha(d, 60)
    lw = 90*PX
    lh = 1*PX
    yy = int(y + size/2)
    for sign in (-1, 1):
        if sign < 0:
            lx = int(cx - size/2 - 14*PX - lw)
        else:
            lx = int(cx + size/2 + 14*PX)
        ln = add_rect(slide, lx, yy, lw, max(lh, 2*PX), color)
        set_alpha(ln, 40)

def top_line(slide):
    add_rect(slide, 0, 0, SLIDE_W, 3*PX, GOLD)

def footer(slide):
    txt(slide, 0, SLIDE_H-30*PX, SLIDE_W, 22*PX, FOOTER,
        font=SANS, size=8, color=MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        spacing=1.5)

def overline(slide, label, color=GOLD):
    txt(slide, 80*PX, 40*PX, 1000*PX, 24*PX, label.upper(),
        font=SANS, size=10, color=color, bold=True, spacing=3)

def clean(s):
    if not s: return ''
    s = re.sub(r'\s+', ' ', s).strip()
    s = re.sub(r'Objetivo:.*$', '', s).strip()
    return s

def split_bullets(s):
    s = clean(s)
    if '●' in s:
        idx = s.find('●')
        intro = s[:idx].strip().rstrip(':')
        parts = [p.strip(' ;.') for p in s[idx:].split('●') if p.strip(' ;.')]
        return intro, parts
    return s, []

# ---------------- LAYOUTS ----------------
def layout_capa(slide, d):
    set_bg(slide, VOID)
    apply_photo(slide, d['num'])
    top_line(slide)
    add_rect(slide, 0, SLIDE_H-3*PX, SLIDE_W, 3*PX, GOLD)
    txt(slide, 0, 120*PX, SLIDE_W, 24*PX, "ESTRUTURA NARRATIVA INSTITUCIONAL — VERSÃO MASTER",
        font=SANS, size=10, color=GOLD, bold=True, align=PP_ALIGN.CENTER, spacing=3)
    txt(slide, 0, 200*PX, SLIDE_W, 90*PX, "JUBIABÁ",
        font=SERIF, size=80, color=OFFWHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        spacing=4)
    txt(slide, 0, 300*PX, SLIDE_W, 50*PX, "A Próxima Fronteira do Luxo no Nordeste",
        font=SERIF, size=30, color=GOLD_L, italic=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    diamond(slide, SLIDE_W/2, 370*PX)
    sub = clean(d.get('subheadline','Um território raro no momento exato da transformação do litoral brasileiro.'))
    txt(slide, 240*PX, 400*PX, SLIDE_W-480*PX, 50*PX, sub,
        font=SANS, size=14, color=CREAM, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)
    n = len(STATS_CAPA)
    cw = 240*PX; gap = 30*PX
    total = n*cw + (n-1)*gap
    sx = int(SLIDE_W/2 - total/2)
    sy = 480*PX
    for i,(val,lab) in enumerate(STATS_CAPA):
        x = sx + i*(cw+gap)
        card = add_rect(slide, x, sy, cw, 90*PX, STONE)
        set_alpha(card, 88)
        txt(slide, x, sy+14*PX, cw, 40*PX, val, font=SERIF, size=30, color=GOLD,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(slide, x, sy+58*PX, cw, 22*PX, lab, font=SANS, size=9, color=MUTED,
            align=PP_ALIGN.CENTER, spacing=2)
    txt(slide, 0, 620*PX, SLIDE_W, 24*PX, "GSWP   ×   KÖNIGSBERGER VANNUCCHI   ×   HAMAM GLOBAL",
        font=SANS, size=11, color=CREAM, align=PP_ALIGN.CENTER, spacing=2)

def layout_capitulo(slide, d, chap_label):
    set_bg(slide, CHARCOAL)
    apply_photo(slide, d['num'])
    top_line(slide)
    overline(slide, chap_label)
    title = d['title']
    txt(slide, 120*PX, 0, SLIDE_W-240*PX, SLIDE_H, title,
        font=SERIF, size=46, color=OFFWHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        line_spacing=1.05)
    diamond(slide, SLIDE_W/2, 470*PX)
    hl = clean(d.get('headline',''))
    if hl:
        txt(slide, 200*PX, 490*PX, SLIDE_W-400*PX, 60*PX, hl,
            font=SERIF, size=20, color=GOLD_L, italic=True, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
    footer(slide)

def layout_final(slide, d):
    set_bg(slide, VOID)
    apply_photo(slide, d['num'])
    top_line(slide)
    add_rect(slide, 0, SLIDE_H-3*PX, SLIDE_W, 3*PX, GOLD)
    txt(slide, 0, 130*PX, SLIDE_W, 24*PX, "JUBIABÁ",
        font=SANS, size=11, color=GOLD, align=PP_ALIGN.CENTER, spacing=4, bold=True)
    title = clean(d.get('headline','') or d['title'])
    txt(slide, 120*PX, 190*PX, SLIDE_W-240*PX, 120*PX, title,
        font=SERIF, size=40, color=GOLD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        line_spacing=1.1)
    diamond(slide, SLIDE_W/2, 330*PX)
    finals = [("R$ 2 bi","VGV"),("250 ha","ÁREA"),("1.070 m","FRENTE-MAR"),("3,5×","RETORNO OTIMISTA")]
    n=len(finals); cw=230*PX; gap=20*PX
    total=n*cw+(n-1)*gap; sx=int(SLIDE_W/2-total/2); sy=380*PX
    for i,(val,lab) in enumerate(finals):
        x=sx+i*(cw+gap)
        card = add_rect(slide, x, sy, cw, 90*PX, STONE)
        set_alpha(card, 88)
        txt(slide, x, sy+14*PX, cw, 40*PX, val, font=SERIF, size=26, color=GOLD,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(slide, x, sy+58*PX, cw, 22*PX, lab, font=SANS, size=8, color=MUTED,
            align=PP_ALIGN.CENTER, spacing=1.5)
    txt(slide, 0, 520*PX, SLIDE_W, 30*PX, "GSWP × KÖNIGSBERGER VANNUCCHI × HAMAM GLOBAL",
        font=SANS, size=12, color=CREAM, align=PP_ALIGN.CENTER, spacing=2)
    txt(slide, 0, 560*PX, SLIDE_W, 24*PX, "Material Confidencial — Apresentação Institucional",
        font=SANS, size=10, color=MUTED, align=PP_ALIGN.CENTER, spacing=1)
    footer(slide)

def layout_tickets(slide, d, chap_label):
    set_bg(slide, VOID)
    apply_photo(slide, d['num'])
    top_line(slide)
    overline(slide, chap_label)
    txt(slide, 80*PX, 80*PX, SLIDE_W-160*PX, 60*PX, d['title'],
        font=SERIF, size=36, color=CREAM, line_spacing=1.0)
    hl = clean(d.get('headline',''))
    if hl:
        txt(slide, 80*PX, 150*PX, SLIDE_W-160*PX, 40*PX, hl,
            font=SERIF, size=18, color=GOLD_L, italic=True)
    n=len(TICKETS); cw=340*PX; gap=30*PX
    total=n*cw+(n-1)*gap; sx=int(SLIDE_W/2-total/2); sy=250*PX
    for i,(name,price,area) in enumerate(TICKETS):
        x=sx+i*(cw+gap)
        star = '★' in name
        card = add_rect(slide, x, sy, cw, 260*PX, STONE)
        if star:
            add_rect(slide, x, sy, cw, 4*PX, GOLD)
        txt(slide, x, sy+30*PX, cw, 30*PX, name, font=SANS, size=13, color=GOLD,
            align=PP_ALIGN.CENTER, bold=True, spacing=2)
        txt(slide, x, sy+90*PX, cw, 50*PX, price, font=SERIF, size=34, color=OFFWHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(slide, x, sy+160*PX, cw, 30*PX, area, font=SANS, size=14, color=CREAM,
            align=PP_ALIGN.CENTER)
        txt(slide, x, sy+200*PX, cw, 24*PX, "ÁREA PRIVATIVA", font=SANS, size=8, color=MUTED,
            align=PP_ALIGN.CENTER, spacing=2)
    footer(slide)

def layout_content(slide, d, chap_label):
    set_bg(slide, VOID)
    photo_kind = apply_photo(slide, d['num'])
    top_line(slide)
    overline(slide, chap_label)
    # se foto lateral, restringe a largura do conteúdo à metade esquerda
    content_w = (SLIDE_W - 160*PX) if photo_kind != 'side' else int(SLIDE_W * 0.52)
    y = 80*PX
    txt(slide, 80*PX, y, content_w, 80*PX, d['title'],
        font=SERIF, size=34, color=CREAM, line_spacing=1.0)
    y += 76*PX
    hl = clean(d.get('headline',''))
    if hl:
        th = 60*PX if len(hl) > 60 else 36*PX
        txt(slide, 80*PX, y, content_w, th, hl,
            font=SERIF, size=20, color=GOLD_L, italic=True, line_spacing=1.15)
        y += th + 14*PX
    diamond(slide, 80*PX+90*PX, y); y += 26*PX
    body = clean(d.get('texto',''))
    dados = d.get('dados','')
    bull = d.get('bullets','')
    intro, items = split_bullets(body)
    extra_intro, extra_items = split_bullets(dados or bull)
    if intro:
        ih = min(120*PX, 30*PX + 18*PX*(len(intro)//70))
        txt(slide, 80*PX, y, content_w, ih, intro,
            font=SANS, size=14, color=CREAM, line_spacing=1.3)
        y += ih
    all_items = items + extra_items
    if all_items:
        if len(all_items) > 4 and photo_kind != 'side':
            half = (len(all_items)+1)//2
            colw = (content_w-40*PX)//2
            add_bullets(slide, 80*PX, y, colw, SLIDE_H-y-50*PX, all_items[:half], size=13)
            add_bullets(slide, 80*PX+colw+40*PX, y, colw, SLIDE_H-y-50*PX, all_items[half:], size=13)
        else:
            add_bullets(slide, 80*PX, y, content_w, SLIDE_H-y-50*PX, all_items, size=14)
    ref = REFS.get(d['num'])
    if ref:
        txt(slide, 80*PX, SLIDE_H-52*PX, content_w, 20*PX, "FONTE: "+ref,
            font=SANS, size=8, color=GOLD, spacing=1)
    footer(slide)

# ---------------- CLASSIFICAÇÃO DE SLIDES ----------------
CHAPTERS = [
    (1,1,"CAPA"),
    (2,20,"BLOCO 01 — A TESE"),
    (21,40,"BLOCO 02 — O TERRITÓRIO"),
    (41,72,"BLOCO 03 — A OPORTUNIDADE"),
    (73,89,"BLOCO 04 — EXECUÇÃO E GOVERNANÇA"),
]
def chap_label(num):
    n=int(num)
    for a,b,lab in CHAPTERS:
        if a<=n<=b: return lab
    return "JUBIABÁ"

CHAPTER_SLIDES = {'20','40','60','72','73','89'}
TICKET_SLIDE = '55'

def main():
    with open(os.path.join(BASE,'_slides_parsed.json'),encoding='utf-8') as f:
        slides = json.load(f)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    photo_slides = []
    for d in slides:
        num = d['num']
        slide = prs.slides.add_slide(blank)
        lbl = chap_label(num)
        if num == '01':
            layout_capa(slide, d)
        elif num == '89':
            layout_final(slide, d)
        elif num == TICKET_SLIDE:
            layout_tickets(slide, d, lbl)
        elif num in CHAPTER_SLIDES:
            layout_capitulo(slide, d, lbl)
        else:
            layout_content(slide, d, lbl)
        if num in SLIDE_PHOTOS:
            photo_slides.append(num)

    out = os.path.join(BASE,'jubiaba-gswp-brand-v2.pptx')
    prs.save(out)
    sz = os.path.getsize(out)
    print(f"SAVED: {out}")
    print(f"SLIDES: {len(prs.slides._sldIdLst)}")
    print(f"SIZE_KB: {sz/1024:.1f}")
    print(f"PIL: {HAS_PIL}")
    print(f"PHOTO_SLIDES ({len(photo_slides)}): {', '.join(photo_slides)}")

if __name__ == '__main__':
    main()
